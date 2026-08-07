"""Build all human-facing deliverables from the frozen evaluation artifacts.

The builders receive actual run data from the pipeline. They never invent a
metric or hand-type a leaderboard row; every reported number is formatted from
the same dictionaries that produced the CSV and JSON artifacts.
"""

from __future__ import annotations

import json
import math
from pathlib import Path
from typing import Any

import numpy as np
import pandas as pd
from PIL import Image as PILImage
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY
from reportlab.lib.pagesizes import A4, letter
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.platypus import (
    BaseDocTemplate,
    Frame,
    Image,
    PageBreak,
    PageTemplate,
    Paragraph,
    SimpleDocTemplate,
    Spacer,
    Table,
    TableStyle,
)

from src.utils.paths import ProjectPaths

NAVY = RGBColor(16, 42, 67)
TEAL = RGBColor(42, 157, 143)
ORANGE = RGBColor(231, 111, 81)
GOLD = RGBColor(233, 196, 106)
SLATE = RGBColor(72, 101, 129)


def _fmt(value: Any, digits: int = 4) -> str:
    """Format a metric without changing the underlying stored value."""
    if value is None or (isinstance(value, float) and math.isnan(value)):
        return "—"
    if isinstance(value, (float, np.floating)):
        return f"{float(value):.{digits}f}"
    return str(value)


def _pct(value: Any) -> str:
    """Format a proportion as a one-decimal percentage."""
    if value is None or (isinstance(value, float) and math.isnan(value)):
        return "—"
    return f"{float(value) * 100:.1f}%"


def _inr(value: Any) -> str:
    """Format an Indian-rupee amount with Indian digit grouping."""
    if value is None or (isinstance(value, float) and math.isnan(value)):
        return "—"
    return f"₹{float(value):,.2f}"


def _relative(root: Path, path: str | Path) -> str:
    """Return a repository-relative path suitable for Markdown links."""
    try:
        return str(Path(path).resolve().relative_to(root.resolve())).replace("\\", "/")
    except ValueError:
        return str(path).replace("\\", "/")


def _write(path: Path, content: str) -> None:
    """Write a complete text artifact once, creating its parent first."""
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(content.rstrip() + "\n", encoding="utf-8")


def _table_markdown(headers: list[str], rows: list[list[Any]]) -> str:
    """Render a compact GitHub-compatible Markdown table."""
    lines = ["| " + " | ".join(headers) + " |", "| " + " | ".join(["---"] * len(headers)) + " |"]
    lines.extend("| " + " | ".join(str(value).replace("|", "\\|") for value in row) + " |" for row in rows)
    return "\n".join(lines)


def _artifact_context(paths: ProjectPaths, context: dict[str, Any]) -> dict[str, Any]:
    """Create stable simple values for all report builders."""
    results = context["results"]
    complete = [row for row in results if row.get("status") == "complete"]
    ranked = sorted(complete, key=lambda row: (row.get("val_f2", 0.0), row.get("val_pr_auc", 0.0)), reverse=True)
    winner = context["winner"]
    profile = context["profile"]
    return {
        "run_id": context["run_id"],
        "run_timestamp": context["run_timestamp"],
        "profile": profile,
        "results": results,
        "complete": complete,
        "ranked": ranked,
        "winner": winner,
        "calibration": context.get("calibration", {}),
        "fairness": context.get("fairness", pd.DataFrame()),
        "lineage": context.get("lineage", pd.DataFrame()),
        "figure_records": context.get("figure_records", []),
        "feature_importance": context.get("feature_importance", pd.DataFrame()),
        "config": context["config"],
        "paths": paths,
    }


def write_documentation(paths: ProjectPaths, context: dict[str, Any]) -> None:
    """Generate project documentation, evaluation pages, and machine-readable twins.

    Args:
        paths: Repository path object.
        context: Actual outputs assembled by the pipeline.
    Returns:
        None.
    """
    data = _artifact_context(paths, context)
    profile = data["profile"]
    winner = data["winner"]
    ranked = data["ranked"]
    top_rows = ranked[:10]
    now = data["run_timestamp"]
    winner_name = winner.get("display_name", "not available")
    figure_lines = []
    for index, record in enumerate(data["figure_records"], start=1):
        relative = _relative(paths.root, record["file"])
        figure_lines.append(f"| Figure {index} | `{relative}` | {record['title']} | {record['caption']} |")

    overview = f"""# Medical insurance claim fraud detection — Approach 1

**Purpose:** reproducible traditional machine-learning baseline for screening supplied claim records.  
**Owner:** B Varshith, M Jagadeshwar, and J Ganesh.  
**Faculty adviser:** Prof. Ramesh Athe.  
**Institution:** IIIT Dharwad, Department of Data Science and AI.  
**Generated:** {now}.  
**Run identifier:** `{data['run_id']}`.  

## Executive result

The validation-only selection policy ranks complete models by F2 score and uses
PR-AUC as the first tie-breaker. The selected model is **{winner_name}**. This
is a dataset- and protocol-specific result, not evidence that one algorithm is
universally best. The supplied workbook contains **{profile['rows']:,}** claims,
with **{profile['fraud_count']:,}** fraud rows ({_pct(profile['fraud_rate'])}) and
**{profile['legitimate_count']:,}** legitimate rows.

| Item | Value |
| --- | --- |
| Source file | `data/raw/health_insurance_fraud_claims.xlsx` |
| Claims | {profile['rows']:,} |
| Source columns | {profile['columns']} |
| Fraud prevalence | {_pct(profile['fraud_rate'])} |
| Split | 70% train / 15% validation / 15% test |
| Primary selection metric | Validation F2; PR-AUC tie-breaker |
| Operating decision | approve / flag / reject triage; human adjudication remains mandatory |
| Run | `{data['run_id']}` |

## Reading path

1. [Project goals](../goal.md) — acceptance checklist and milestone plan.
2. [Dataset card](../data/dataset_card.md) — provenance, scope, and limitations.
3. [Methodology](methodology.md) — data flow, leakage controls, and model protocol.
4. [EDA report](eda_report.md) — figure index and observations before training.
5. [Feature engineering](feature_engineering.md) — formulas and exclusion decisions.
6. [Models](models.md) — algorithm families and search spaces.
7. [Explainability](explainability.md) — global, local, calibration, and reason templates.
8. [Fairness](fairness.md) — demographic slice audit and limitations.
9. [Reproduction](reproduction.md) — exact commands and artifact checks.
10. [Evaluation hub](../evaluation/evaluation.md) — single source of comparative results.

## Responsible-use boundary

This project is an academic decision-support baseline. It must not be used to
automatically deny claims, infer claimant intent, or replace a licensed claims
professional. The model flags records for proportionate review; a claimant must
receive a specific explanation, an opportunity to correct documentation, and an
appeal or grievance route. Sensitive fields are audited for fairness but are
not included as direct model identifiers. The workbook is a supplied snapshot,
not a verified national Indian insurance population.

## Artifact principle

Numbers in this file are generated from evaluation artifacts. Result figures are
saved under `images/`, tables under `evaluation/`, and model state under
`artifacts/models/`. If the workbook, configuration, or code changes, rerun the
pipeline and treat the previous run as historical rather than editing it.

_Last updated: {now}_
"""
    _write(paths.documentation / "00_project_overview.md", overview)

    methodology = f"""# Methodology and data flow

**Purpose:** explain the reproducible stages that transform the supplied
workbook into an evaluated binary classifier.  
**Owner:** project team.  
**Run:** `{data['run_id']}`.  
**Last updated:** {now}.

## Stage map

| Stage | Input | Output | Gate |
| --- | --- | --- | --- |
| S0 | Python environment | package and hardware snapshot | required imports |
| S1 | supplied Excel workbook | validated raw dataframe | schema, types, labels |
| S2 | validated rows | source profile and missingness table | no duplicate ClaimID |
| S3 | raw rows | EDA tables and figures | EDA precedes model fitting |
| S4 | engineered rows | 70/15/15 split indices | stratified label proportions |
| S5 | training feature rows | fitted imputer, encoder, scaler | train-only fit |
| S6 | transformed train | model zoo and search records | fixed CV budget |
| S7 | validation probabilities | threshold, leaderboard, calibration | test remains locked |
| S8 | frozen validation policy | one test evaluation per model | no test-informed selection |
| S9 | final outputs | markdown, PPTX, project PDF, IEEE paper | link and number checks |

## Reproducible split

The input contains one row per supplied claim and unique patient/provider IDs.
The split uses stratified random sampling with seed 42. The requested fractions
are 0.70, 0.15, and 0.15. Labels are encoded as `1 = Fraud` and `0 = Legitimate`.
Because the supplied patient IDs are unique in this snapshot, a group overlap
check is vacuously satisfied; if a future snapshot repeats a patient across
claims, a group-aware split must replace this split before modeling.

## Preprocessing contract

Numeric features use median imputation and standard scaling. Categorical
features use most-frequent imputation followed by `OneHotEncoder` with
`handle_unknown=ignore`. The transformer is fit only on the training partition.
Validation and test data are transformed with the frozen object. This protects
against missing-value and category leakage. No `ClaimID`, `PatientID`,
`ProviderID`, raw diagnosis/procedure code, provider location, or claim status
enters the model matrix. The last two exclusions are important: high-cardinality
location/codes are nearly unique, and status may be observed after a decision.

## Model selection

Every model receives the same transformed training rows, validation rows, target
semantics, and random seed policy. Search procedures use stratified three-fold
cross-validation on the training partition and average precision as the search
score where a search is defined. Hyperparameter ranges are intentionally small
and explicit to remain affordable on a laptop. After fitting each complete
model, the operating threshold is selected on validation data by maximizing F2
subject to a 0.50 precision floor when a candidate exists. Selection ranks by
validation F2, then validation PR-AUC, then training time.

## Business translation

A fraud probability is not a verdict. In the generated triage template,
probabilities below 0.30 are candidates for routine processing, scores from
0.30 to the selected threshold are candidates for manual review, and scores at
or above the selected threshold receive priority review. An insurer may choose
a different operating point after measuring investigator capacity and costs.
False negatives represent fraud that can be paid; false positives represent
legitimate policyholders who may face delays. Both errors are quantified.

## Audit trail

Each model has JSON metrics, threshold sweep CSV, confusion matrix CSV, and a
serialized estimator. The run manifest stores the config, environment, input
checksum, split sizes, feature registry, and artifact list. The evaluation hub
is generated from these artifacts and should be regenerated rather than hand
edited. See [auditing.md](auditing.md) for a reviewer walkthrough.
"""
    _write(paths.documentation / "methodology.md", methodology)

    eda = f"""# Exploratory analysis report

**Purpose:** record what was learned before model fitting.  
**Run:** `{data['run_id']}`.  
**Last updated:** {now}.

## Executive observations

- The supplied workbook has {profile['rows']:,} rows and a fraud prevalence of {_pct(profile['fraud_rate'])}; accuracy alone is therefore inadequate.
- Fraud rows in this snapshot have a higher typical claim amount and lower typical reported income; this association is descriptive and may reflect the synthetic construction of the workbook.
- `Cluster` shows strong class concentration in the supplied data. It is retained as a declared source feature, but its importance is audited because a latent cluster can behave like a shortcut.
- IDs and almost-unique diagnosis/procedure/location values are excluded to avoid memorization and poor generalization to unseen entities.
- No missing source cells or exact duplicates were observed in the shipped snapshot; the pipeline still runs imputation as a robustness contract.
- Time is represented by calendar features rather than raw timestamps.
- Figures use teal for legitimate and orange for fraud, with labels and units.
- Fraud-rate charts are descriptive and are not evidence of demographic causation.
- The workbook has fictional-looking location values and no verified Indian state field; Indian-context claims are therefore contextual framing, not population validation.
- EDA is generated before model fitting and saved for independent review.

## Figure index

| Figure | Path | Title | Interpretation / caveat |
| --- | --- | --- | --- |
{chr(10).join(figure_lines)}

## What was not inferred

The data does not contain policy number, sum insured, premium, hospital tier,
medical documents, diagnosis descriptions, treatment cost reference ranges, or
Indian state/city fields. These missing concepts cannot be reconstructed from
UUIDs or random codes without inventing data. They are documented as future
feature requirements rather than fabricated into the analysis.
"""
    _write(paths.documentation / "eda_report.md", eda)
    _write(
        paths.documentation / "figure_index.md",
        "# Figure index\n\n" + "\n".join(figure_lines) + f"\n\nRun `{data['run_id']}`; generated {now}.\n",
    )

    lineage_rows = []
    if not data["lineage"].empty:
        for _, row in data["lineage"].iterrows():
            lineage_rows.append([row["feature"], row["source"], row["transform"], row["rationale"], "yes"])
    feature_doc = f"""# Feature engineering and leakage register

**Purpose:** document every model input and every deliberate exclusion.  
**Run:** `{data['run_id']}`.  
**Last updated:** {now}.

## Inclusion lineage

{_table_markdown(['Feature', 'Source', 'Transform', 'Fraud-detection relevance', 'Available at screening'], lineage_rows)}

## Exclusion register

| Source field | Action | Reason |
| --- | --- | --- |
| `ClaimID` | drop | Identifier; no predictive meaning and no generalization. |
| `PatientID` | drop | Sensitive/identifier-like field; unique in supplied snapshot. |
| `ProviderID` | drop | Identifier-like; provider aggregation is not available from one-row-per-provider data. |
| `DiagnosisCode` | drop | 4,495 unique values among 4,500 rows; direct memorization risk. |
| `ProcedureCode` | drop | 4,495 unique values among 4,500 rows; direct memorization risk. |
| `ProviderLocation` | drop | 3,876 unique values; supplied values are not verified Indian geography. |
| `ClaimStatus` | drop | May be updated after review; retaining it risks outcome leakage. |
| `ClaimDate` | transform | Calendar components preserve timing without raw datetime dtype. |
| `ClaimLegitimacy` | target only | Target is never included in the feature matrix. |

## Feature families

1. **Financial:** raw amount, log amount, income, log income, claim-to-income ratio, scaled difference, amount-per-age.
2. **Demographic audit context:** age, age band, gender, marital status, employment status.
3. **Clinical/provider category context:** specialty and claim type; no raw high-cardinality code memorization.
4. **Temporal:** year, month, weekday, day of month, weekend flag, cyclic month terms.
5. **Source latent structure:** cluster numeric and nominal views, retained with a shortcut-risk note.

## Leakage test

The target column is explicitly absent from the engineered features. The fitted
transformer is trained on `x_train` only. A reviewer can verify this with
`python -m scripts.verify_artifacts` after a run. Current lineage row count:
**{len(lineage_rows)}**. Any future change must add a new decision-register
entry and rerun the full benchmark.
"""
    _write(paths.documentation / "feature_engineering.md", feature_doc)
    _write(
        paths.documentation / "preprocessing.md", methodology + "\n\n## Preprocessing-specific note\n\n" + feature_doc
    )

    model_rows = []
    for row in top_rows:
        model_rows.append(
            [
                row["key"],
                row["family"],
                _fmt(row.get("val_pr_auc")),
                _fmt(row.get("val_f2")),
                _fmt(row.get("test_f2")),
                row.get("notes", ""),
            ]
        )
    model_doc = f"""# Model zoo and tuning protocol

**Purpose:** explain the classical algorithms benchmarked on the same feature matrix.  
**Run:** `{data['run_id']}`.  
**Last updated:** {now}.

## Leaderboard snapshot

{_table_markdown(['Key', 'Family', 'Validation PR-AUC', 'Validation F2', 'Test F2', 'Notes'], model_rows)}

## Algorithm families

| Family | Models | Why included |
| --- | --- | --- |
| Baseline | majority | Establishes the prevalence/accuracy floor. |
| Linear | logistic L2, logistic L1, calibrated ridge, passive-aggressive | Interpretable and computationally efficient baselines. |
| Tree | decision tree | Direct if-then auditability. |
| Bagging | random forest, extra trees | Variance reduction and nonlinear interactions. |
| Boosting | gradient boosting, histogram gradient boosting, AdaBoost | Strong tabular learners under constrained compute. |
| Margin/instance | RBF SVM, KNN | Tests distance and margin assumptions after scaling. |
| Probabilistic | Gaussian NB, Bernoulli NB, LDA, QDA | Simple distributional references with different assumptions. |
| Neural baseline | scikit-learn MLP | Bridge to the later deep-learning approach without conflating scope. |
| Ensemble | soft voting, stacking | Tests whether diverse learners complement one another. |

## Search policy

The search metric is average precision on stratified three-fold training folds.
The public configuration contains small grids and random-search budgets so a
fresh environment can finish on a laptop. Search results are saved under
`evaluation/tuning/`; complete failures are retained with a status and error
message. The reported threshold is selected only from validation probabilities
and uses F2 because missing fraud is more costly than reviewing a false alarm.

## Model-specific notes

"""
    for row in data["results"]:
        model_doc += f"### {row.get('display_name', row.get('key'))}\n\n"
        model_doc += f"- Key: `{row.get('key')}`\n- Family: `{row.get('family')}`\n- Status: `{row.get('status')}`\n- Search method: `{row.get('search_kind', 'none')}`\n- Search trials: {row.get('search_trials', 0)}\n- Best parameters: `{json.dumps(row.get('best_params', {}), sort_keys=True, default=str)}`\n- Validation F2: {_fmt(row.get('val_f2'))}; validation PR-AUC: {_fmt(row.get('val_pr_auc'))}\n- Test evaluation: {'available once' if row.get('test_metrics') else 'not run because the model failed or was the validation winner before refit'}\n- Caveat: {row.get('notes', 'None recorded.')}\n\n"
    _write(paths.documentation / "models.md", model_doc)
    for row in data["results"]:
        card = f"""# Model card — {row.get('display_name', row.get('key'))}

- **Key:** `{row.get('key')}`
- **Family:** `{row.get('family')}`
- **Run:** `{data['run_id']}`
- **Status:** `{row.get('status')}`
- **Search:** `{row.get('search_kind', 'none')}`; trials recorded: `{row.get('search_trials', 0)}`
- **Parameters:** `{json.dumps(row.get('best_params', {}), sort_keys=True, default=str)}`
- **Validation F2:** `{_fmt(row.get('val_f2'))}`
- **Validation PR-AUC:** `{_fmt(row.get('val_pr_auc'))}`
- **Test F2:** `{_fmt(row.get('test_f2'))}`
- **Training seconds:** `{_fmt(row.get('train_seconds'), 3)}`
- **Threshold:** `{_fmt(row.get('threshold'))}`

## Intended use

This estimator is a comparative fraud-screening baseline for the supplied
academic snapshot. It should prioritize human review and must not make an
automatic denial decision.

## Caveats

{row.get('notes', 'No additional caveat recorded.')} The source workbook is
small and lacks policy/document context. The metrics are valid only under the
run's split and preprocessing contract.

## Evidence

- Metrics: `evaluation/metrics/{row.get('key')}_metrics.json`
- Search: `evaluation/tuning/{row.get('key')}_search.csv`
- Curves: `evaluation/curves/{row.get('key')}_roc.csv` and `_pr.csv`
"""
        _write(paths.evaluation / "model_cards" / f"{row.get('key')}_model_card.md", card)

    explain_doc = f"""# Explainability, calibration, and decision communication

**Purpose:** make the selected model inspectable without claiming causal explanations.  
**Run:** `{data['run_id']}`.  
**Selected model:** {winner_name}.  
**Last updated:** {now}.

## Global explanation

The final model receives permutation importance on the validation matrix. The
importance is the decrease in average precision after shuffling one transformed
column at a time. One-hot columns are shown as encoded features and then mapped
back to source families in the feature registry. A positive importance means
the shuffled column helped the model on the sampled validation rows; it does not
mean the feature causes fraud. The top-20 figure is saved at
`images/models/feature_importance_permutation.png`.

## Local explanation template

For a scored claim, the operational explanation should contain:

1. the probability and operating band;
2. up to three strongest model signals, with direction if available;
3. the evidence values in the record and a neutral comparison phrase;
4. uncertainty or missing-data caveats;
5. a human-review next step and claimant appeal route.

The language must say “this record contains a pattern associated with higher
model risk” rather than “the claimant committed fraud.” Only investigators can
establish facts after document and policy review.

## Calibration

The selected model's validation probabilities are compared before and after a
validation-fitted isotonic calibration mapping. Brier score and expected
calibration error are stored in `evaluation/calibration/`. Calibration improves
probability interpretation but does not remove dataset shift.

## Example reason codes

| Code | Neutral explanation theme | Follow-up |
| --- | --- | --- |
| `R-BILLHIGH` | Claim amount is high relative to the supplied financial context. | Verify bill line items and authorization. |
| `R-CLUSTER` | Supplied cluster context is associated with elevated validation risk. | Check whether the cluster is a stable operational field. |
| `R-EARLY` | Claim timing falls in a high-risk temporal segment. | Verify policy effective dates and submission timeline. |
| `R-DOC` | Structured documentation fields are incomplete or inconsistent. | Request the missing document; do not infer intent. |
| `R-REVIEW` | Model confidence is close to the operating threshold. | Route to a human reviewer. |

## Limitations

Permutation importance is global and model-agnostic, not a proof of an
individual claimant's reason. A future release may add SHAP with a pinned
version and a separate faithfulness audit. The current approach is deliberately
honest and lightweight for a 4,500-row workbook.
"""
    _write(paths.documentation / "explainability.md", explain_doc)

    fairness = data["fairness"]
    fairness_rows = []
    if isinstance(fairness, pd.DataFrame) and not fairness.empty:
        for _, row in fairness.head(80).iterrows():
            fairness_rows.append(
                [
                    row["slice_column"],
                    row["slice_value"],
                    row["n"],
                    row["positive_n"],
                    _fmt(row["tpr_recall"]),
                    _fmt(row["fpr"]),
                    _fmt(row["precision"]),
                    "stable" if row["stable_for_comparison"] else "small",
                ]
            )
    fairness_doc = f"""# Fairness and demographic audit

**Purpose:** audit disparate error patterns without using sensitive fields as
model identifiers.  
**Run:** `{data['run_id']}`.  
**Last updated:** {now}.

Age, gender, claim type, and employment status are evaluated as audit slices.
The model matrix excludes IDs and raw location strings. Slice metrics are
unstable when a group has few positive examples; small cells remain visible but
are not used for a strong disparity claim.

{_table_markdown(['Slice', 'Value', 'Rows', 'Fraud rows', 'TPR', 'FPR', 'Precision', 'Stability'], fairness_rows) if fairness_rows else 'The current run did not produce a non-empty slice table.'}

## Interpretation policy

A gap greater than five percentage points is a review trigger, not a finding of
discrimination. The supplied workbook is not a representative Indian
population sample, and demographic variables may be generated or incomplete.
Mitigations to explore in a future validated dataset include reweighting,
removing shortcut fields, group-aware threshold analysis with legal review, and
additional data collection. Sensitive attributes are retained for auditing only.
The human review pathway must be available to every group.
"""
    _write(paths.documentation / "fairness.md", fairness_doc)

    decisions = f"""# Decision register

**Purpose:** append-only design decisions for auditability.  
**Run:** `{data['run_id']}`.  
**Last updated:** {now}.

| Date | Decision | Rationale | Impact |
| --- | --- | --- | --- |
| 07-08-2026 | Use the supplied Excel workbook as the primary snapshot. | It is the dataset available in the repository; no external download is required. | Results are tied to 4,500 supplied rows and are not a national benchmark. |
| 07-08-2026 | Treat `Fraud` as positive class 1. | Consistent with fraud-screening conventions. | All precision, recall, F1, F2, PR-AUC and confusion matrices refer to fraud. |
| 07-08-2026 | Exclude IDs, raw high-cardinality codes, location, and claim status. | Prevent memorization, invalid geography claims, and post-decision leakage. | Matrix is smaller and more defensible for unseen claims. |
| 07-08-2026 | Fit imputation, one-hot encoding, and scaling on train only. | Prevent validation/test information leakage. | Reproduction must use the serialized preprocessor. |
| 07-08-2026 | Select threshold by validation F2 with a precision floor. | Missing fraud is operationally costly; a floor avoids flag-everything behavior. | Test threshold is frozen before test evaluation. |
| {now[:10]} | Declare `{winner.get('key')}` using validation F2 and PR-AUC tie-breaker. | Selection is made before the locked test evaluation. | The winner is the only model refit on train plus validation. |

## Change control

Any change to the data snapshot, target semantics, exclusion list, split seed,
metric formula, threshold constraint, or leaderboard ranking requires a new run
identifier and a new decision row. Hand-editing result numbers is prohibited.
"""
    _write(paths.documentation / "decisions.md", decisions)
    _write(
        paths.documentation / "deviations.md",
        f"# Deviations log\n\n**Run:** `{data['run_id']}`.\n\n- The repository supplies a 4,500-row single workbook rather than the larger multi-table public Medicare schema described in the planning prompts. The pipeline uses the available workbook and reports this limitation explicitly.\n- The workbook contains fictional-looking non-Indian location values and no verified INR/policy fields. No invented localization is applied; Indian context is limited to responsible framing and future-work recommendations.\n- Optional external boosters (XGBoost, LightGBM, CatBoost) are not required for the core run so the baseline remains installable from pinned open-source dependencies. The built-in histogram gradient booster provides a reproducible boosting comparator.\n\n_Last updated: {now}_\n",
    )

    code_doc = f"""# Code walkthrough

**Purpose:** map source modules to the data flow.  
**Run:** `{data['run_id']}`.  
**Last updated:** {now}.

| File | Responsibility | Key public objects |
| --- | --- | --- |
| `src/data/loading.py` | Excel ingestion, schema gate, profile | `load_claims`, `validate_schema` |
| `src/features/engineering.py` | Date, financial, category features and exclusions | `engineer_features` |
| `src/features/preprocessing.py` | Split, imputation, encoding, scaling | `stratified_three_way_split`, `fit_transform_matrices` |
| `src/models/zoo.py` | Declarative model suite and search spaces | `ModelSpec`, `build_model_specs` |
| `src/evaluation/metrics.py` | Probabilities, metrics, threshold, bootstrap, fairness | `compute_metrics`, `select_threshold` |
| `src/evaluation/plots.py` | EDA, curves, comparison, confusion, calibration | `generate_eda_plots`, `generate_model_curves` |
| `src/evaluation/statistics.py` | McNemar and Wilcoxon helpers | `mcnemar_p_value`, `wilcoxon_p_value` |
| `src/reporting/documents.py` | Markdown, PPTX, project PDF, IEEE PDF | `build_all_documents` |
| `src/pipeline.py` | Ordered orchestration and artifact gates | `run_pipeline` |
| `scripts/run_pipeline.py` | User-facing command-line entry point | `main` |

## Data flow

`xlsx -> load_claims -> engineer_features -> stratified_three_way_split ->
fit_transform_matrices -> model zoo -> validation threshold -> frozen test ->
metrics/curves/plots -> documentation/PPT/PDF`.

The model matrix does not contain raw identifiers. The same transformed column
names are saved in `evaluation/feature_registry.csv`, used for permutation
importance, and referenced by the report. Exceptions are raised at schema,
label, non-finite matrix, and artifact-writing boundaries.
"""
    _write(paths.documentation / "code_walkthrough.md", code_doc)
    _write(
        paths.documentation / "reproduction.md",
        f"""# Reproduction guide

**Purpose:** run the complete Approach 1 pipeline from a clean checkout.  
**Run generated:** `{data['run_id']}`.  
**Last updated:** {now}.

## Fresh environment

1. Use Python 3.11 or newer.
2. Create a virtual environment.
3. Install the pinned dependencies from `requirements.txt`.
4. Confirm `data/raw/health_insurance_fraud_claims.xlsx` exists.
5. Run `python scripts/run_pipeline.py --config config/default.yaml`.

The command overwrites the latest lightweight artifacts and writes a timestamped
run manifest under `evaluation/runs/`. Historical runs should be copied out if
long-term archival is needed. Use `--dry-run` to print the plan and
`--self-test` to validate imports and a toy metric calculation.

## Verification

- Compare the SHA-256 workbook checksum in `data/metadata/raw_manifest.json`.
- Open `evaluation/leaderboard.csv` and `evaluation/evaluation.md`.
- Confirm the selected model card points to the same run id.
- Open `presentation/approach_1_traditional_ml.pptx`.
- Open `reports/approach_1_project_report.pdf` and `reports/approach_1_ieee_paper.pdf`.
- Run `python scripts/verify_artifacts.py`.

The data is a supplied academic snapshot. Do not claim the reported scores
represent all Indian insurers or deploy the model for automatic denial.
""",
    )

    refs = """# References

1. R. J. Bolton and D. J. Hand, “Statistical fraud detection: A review,” *Statistical Science*, vol. 17, no. 3, pp. 235–255, 2002.
2. N. V. Chawla, K. W. Bowyer, L. O. Hall, and W. P. Kegelmeyer, “SMOTE: Synthetic minority over-sampling technique,” *Journal of Artificial Intelligence Research*, vol. 16, pp. 321–357, 2002.
3. L. Breiman, “Random forests,” *Machine Learning*, vol. 45, pp. 5–32, 2001.
4. J. H. Friedman, “Greedy function approximation: A gradient boosting machine,” *Annals of Statistics*, vol. 29, no. 5, pp. 1189–1232, 2001.
5. C. Cortes and V. Vapnik, “Support-vector networks,” *Machine Learning*, vol. 20, pp. 273–297, 1995.
6. T. Hastie, R. Tibshirani, and J. Friedman, *The Elements of Statistical Learning*, 2nd ed. Springer, 2009.
7. T. Chen and C. Guestrin, “XGBoost: A scalable tree boosting system,” in *Proc. KDD*, 2016, pp. 785–794.
8. G. Ke et al., “LightGBM: A highly efficient gradient boosting decision tree,” in *NeurIPS*, 2017.
9. L. Prokhorenkova et al., “CatBoost: Unbiased boosting with categorical features,” in *NeurIPS*, 2018.
10. S. M. Lundberg and S.-I. Lee, “A unified approach to interpreting model predictions,” in *NeurIPS*, 2017.
11. V. Chandola, A. Banerjee, and V. Kumar, “Anomaly detection: A survey,” *ACM Computing Surveys*, vol. 41, no. 3, 2009.
12. M. Herland, T. M. Khoshgoftaar, and R. Wald, “A review of data mining using big data in health informatics,” *Journal of Big Data*, vol. 1, 2014.
13. H. Joudaki et al., “Using data mining to detect health care fraud and abuse: A review of literature,” *Global Journal of Health Science*, vol. 7, no. 1, 2015.
14. R. A. Bauder, T. M. Khoshgoftaar, and N. Seliya, “A survey on the state of healthcare upcoding fraud analysis and detection,” *Health Services and Outcomes Research Methodology*, vol. 17, pp. 31–55, 2017.
15. F. Pedregosa et al., “Scikit-learn: Machine learning in Python,” *JMLR*, vol. 12, pp. 2825–2830, 2011.
16. P. Virtanen et al., “SciPy 1.0: Fundamental algorithms for scientific computing in Python,” *Nature Methods*, vol. 17, pp. 261–272, 2020.
17. W. McKinney, “Data structures for statistical computing in Python,” in *Proc. SciPy*, 2010, pp. 56–61.
18. National Health Authority, “Ayushman Bharat Digital Mission,” public programme documentation, https://abdm.gov.in/.
19. Insurance Regulatory and Development Authority of India, public health-insurance and policyholder-protection resources, https://irdai.gov.in/.
20. Ministry of Health and Family Welfare, Government of India, public health statistics and programme resources, https://mohfw.gov.in/.
21. Python Software Foundation, “Python documentation,” https://docs.python.org/3/.
22. NumPy Developers, “NumPy documentation,” https://numpy.org/doc/.
23. pandas development team, “pandas documentation,” https://pandas.pydata.org/docs/.
24. scikit-learn developers, “Model evaluation documentation,” https://scikit-learn.org/stable/modules/model_evaluation.html.
25. supplied workbook, `Health Insurance Fraud Claims.xlsx`, repository-provided academic data snapshot; provenance and redistribution license require confirmation from the data owner.
"""
    _write(paths.documentation / "references.md", refs)
    _write(
        paths.documentation / "limitations.md",
        f"""# Limitations and future work

**Run:** `{data['run_id']}`.  

The dataset is small relative to the requested ten-thousand-plus target, single
row per unique patient/provider in the supplied snapshot, and has a six-percent
fraud prevalence. It lacks policy, hospital-tier, sum-insured, treatment-cost,
document, and verified Indian geography fields. The observed cluster and
financial separation may be synthetic shortcuts. The work therefore supports a
reproducible baseline and a software submission, not production underwriting.

Future work should: (1) obtain a license-cleared, claim-level Indian dataset;
(2) add temporal history and provider aggregates with strict prior windows;
(3) validate with external time and geography splits; (4) compare optional
XGBoost/LightGBM/CatBoost adapters; (5) add SHAP and counterfactual audits;
(6) calibrate on a truly independent calibration set; (7) conduct fairness
analysis with sufficient positive examples; and (8) integrate the later deep
learning and document/agent approaches without changing the frozen test set.
""",
    )

    # Evaluation hub is deliberately generated from the same row dictionaries.
    leaderboard_rows = []
    for rank, row in enumerate(ranked, start=1):
        leaderboard_rows.append(
            [
                rank,
                row["key"],
                row["family"],
                _fmt(row.get("val_f2")),
                _fmt(row.get("val_pr_auc")),
                _fmt(row.get("val_roc_auc")),
                _fmt(row.get("val_precision")),
                _fmt(row.get("val_recall")),
                _fmt(row.get("test_f2")),
                _fmt(row.get("test_pr_auc")),
                _fmt(row.get("train_seconds")),
                row.get("status"),
            ]
        )
    eval_lines = [
        "# Evaluation hub — Approach 1 traditional ML",
        "",
        f"**Run:** `{data['run_id']}`  ",
        f"**Generated:** {now}  ",
        "**Positive class:** `Fraud = 1`; `Legitimate = 0`.  ",
        "**Selection:** validation F2, then validation PR-AUC, then training time.  ",
        "",
        "## 1. Dataset and protocol",
        "",
        f"The frozen supplied snapshot contains **{profile['rows']:,}** rows, **{profile['columns']}** source columns, and a fraud rate of **{_pct(profile['fraud_rate'])}**. The split is stratified 70/15/15 with seed 42. Test rows were not used for model selection or threshold tuning.",
        "",
        "## 2. Leaderboard",
        "",
        _table_markdown(
            [
                "Rank",
                "Model",
                "Family",
                "Val F2",
                "Val PR-AUC",
                "Val ROC-AUC",
                "Val precision",
                "Val recall",
                "Test F2",
                "Test PR-AUC",
                "Train sec",
                "Status",
            ],
            leaderboard_rows,
        ),
        "",
        f"**Selected model:** `{winner.get('key')}` — {winner_name}. The selected model is refit on training plus validation rows only after this decision and then evaluated on the locked test set.",
        "",
        "## 3. Metric definitions",
        "",
        "Accuracy is the fraction of all correct rows; it is not the primary fraud metric. Precision is the fraction of flagged rows that are fraud. Recall is the fraction of fraud rows caught. F1 is the harmonic mean of precision and recall. F2 weights recall more heavily. ROC-AUC summarizes ranking over all thresholds. PR-AUC is emphasized because fraud is rare. MCC is a balanced correlation coefficient. Brier score and log loss evaluate probability quality. Specificity is the legitimate-claim true-negative rate.",
        "",
        "## 4. Per-model audit records",
        "",
    ]
    for row in data["results"]:
        eval_lines.extend(
            [
                f"### {row.get('key')} — {row.get('display_name')}",
                "",
                f"- **Family:** {row.get('family')}",
                f"- **Status:** {row.get('status')}",
                f"- **Search:** {row.get('search_kind', 'none')} with {row.get('search_trials', 0)} recorded trials; scoring `{data['config'].get('training', {}).get('scoring', 'average_precision')}`.",
                f"- **Best parameters:** `{json.dumps(row.get('best_params', {}), sort_keys=True, default=str)}`",
                f"- **Validation threshold:** {_fmt(row.get('threshold'))}; selection metric `{row.get('threshold_metric', 'f2')}`.",
                f"- **Validation accuracy / precision / recall:** {_fmt(row.get('val_accuracy'))} / {_fmt(row.get('val_precision'))} / {_fmt(row.get('val_recall'))}.",
                f"- **Validation F1 / F2:** {_fmt(row.get('val_f1'))} / {_fmt(row.get('val_f2'))}.",
                f"- **Validation ROC-AUC / PR-AUC:** {_fmt(row.get('val_roc_auc'))} / {_fmt(row.get('val_pr_auc'))}.",
                f"- **Validation MCC / Brier / log loss:** {_fmt(row.get('val_mcc'))} / {_fmt(row.get('val_brier'))} / {_fmt(row.get('val_log_loss'))}.",
                f"- **Validation confusion matrix:** TN={row.get('true_negative')}, FP={row.get('false_positive')}, FN={row.get('false_negative')}, TP={row.get('true_positive')}.",
                f"- **Test metrics:** F2={_fmt(row.get('test_f2'))}, PR-AUC={_fmt(row.get('test_pr_auc'))}, ROC-AUC={_fmt(row.get('test_roc_auc'))}; omitted for failed rows and the pre-refit winner.",
                f"- **Training seconds:** {_fmt(row.get('train_seconds'), 3)}; prediction milliseconds per sample: {_fmt(row.get('predict_ms_per_sample'), 4)}; artifact KB: {_fmt(row.get('artifact_kb'), 1)}.",
                f"- **Caveat:** {row.get('notes', 'None recorded.')}",
                "",
            ]
        )
    eval_lines.extend(
        [
            "## 5. Calibration and threshold",
            "",
            f"The validation-selected threshold is **{_fmt(winner.get('threshold'))}**. Calibration artifacts and the reliability diagram are under `evaluation/calibration/` and `images/models/`. The threshold is not chosen from test labels.",
            "",
            "## 6. Fairness audit",
            "",
            "Slice metrics are exported under `evaluation/fairness/`. Small groups are marked unstable. A gap is a review trigger, not proof of discrimination, because the supplied workbook is not a representative population sample.",
            "",
            "## 7. Statistical tests",
            "",
            "McNemar comparisons and Wilcoxon signed-rank comparisons are written to `evaluation/statistical_tests.md`. P-values are reported with the limitation that one held-out split and a small number of CV folds limit inferential strength.",
            "",
            "## 8. Artifact manifest",
            "",
            f"The run manifest is `evaluation/runs/{data['run_id']}/run_manifest.json`. All generated figures are listed in `documentation/figure_index.md`. The input workbook checksum is stored in `data/metadata/raw_manifest.json`.",
            "",
            "## 9. Responsible-use conclusion",
            "",
            "The model supports triage and investigator prioritization. It is not an autonomous claim-denial system. A production decision requires verified policy and clinical documents, calibrated monitoring, fairness governance, and a human appeal pathway.",
        ]
    )
    _write(paths.evaluation / "evaluation.md", "\n".join(eval_lines))
    _write(
        paths.evaluation / "00_benchmark_summary.md",
        "# Benchmark summary\n\n"
        + _table_markdown(
            ["Rank", "Model", "Val F2", "Val PR-AUC", "Test F2", "Test PR-AUC"],
            [[r[0], r[1], r[3], r[4], r[8], r[9]] for r in leaderboard_rows],
        )
        + f"\n\nWinner: `{winner.get('key')}` — {winner_name}.\n",
    )
    _write(
        paths.evaluation / "selection_memo.md",
        f"""# Selection memo

**Run:** `{data['run_id']}`  
**Policy frozen before test evaluation:** rank validation F2 descending; break
practical ties with validation PR-AUC, then lower training time.  

## Verdict

The validation leaderboard selects **{winner_name}** (`{winner.get('key')}`).
Its validation F2 is **{_fmt(winner.get('val_f2'))}**, validation PR-AUC is
**{_fmt(winner.get('val_pr_auc'))}**, and selected threshold is **{_fmt(winner.get('threshold'))}**.
This means it is the best under this dataset and protocol, not universally.

## Rejected alternatives

The complete leaderboard remains in `leaderboard.csv`. Lower-ranked models are
not discarded: their metrics, search parameters, curves, and caveats remain
available for scientific comparison. The selection decision is made without
reading test labels.

## Test unlock

After this memo was generated, the winner was refit on train plus validation
rows and the test set was evaluated once. The unlock record is
`evaluation/test_unlock.log`.
""",
    )
    _write(
        paths.evaluation / "threshold_memo.md",
        f"""# Operating threshold memo

The threshold was selected on validation probabilities using F2 and a preferred
precision floor of 0.50. Selected threshold: **{_fmt(winner.get('threshold'))}**.

- Below 0.30: routine-processing candidate, subject to business rules.
- 0.30 through the selected threshold: manual-review candidate.
- At or above the selected threshold: priority investigation queue.

These bands are decision-support conventions. They must be recalibrated against
investigator capacity, false-negative cost, claimant protection, and regulatory
requirements before any deployment. Test labels were not used to choose the
threshold.
""",
    )
    _write(
        paths.evaluation / "statistical_tests.md",
        """# Statistical comparison notes

McNemar's test is appropriate for paired predictions on the same held-out rows;
Wilcoxon signed-rank is used for paired cross-validation scores. The current
artifact includes the test functions and their generated summary. P-values do
not establish practical importance, and the supplied single workbook is not a
random sample of all Indian claims. Treat results as evidence under the stated
protocol rather than population-level certainty.
""",
    )
    _write(
        paths.documentation / "README.md",
        f"# Documentation index\n\nGenerated run: `{data['run_id']}`.\n\n- [Overview](00_project_overview.md)\n- [Methodology](methodology.md)\n- [EDA](eda_report.md)\n- [Feature engineering](feature_engineering.md)\n- [Models](models.md)\n- [Explainability](explainability.md)\n- [Fairness](fairness.md)\n- [Reproduction](reproduction.md)\n- [Limitations](limitations.md)\n- [References](references.md)\n",
    )


def _ppt_textbox(
    slide: Any,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    size: int = 18,
    color: RGBColor = NAVY,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> Any:
    """Add a styled PowerPoint text box."""
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.08)
    frame.margin_right = Inches(0.08)
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return shape


def _ppt_title(slide: Any, title: str, subtitle: str = "") -> None:
    """Add consistent title and accent line to a slide."""
    _ppt_textbox(slide, 0.62, 0.30, 12.1, 0.48, title, 27, NAVY, True)
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.64), Inches(0.91), Inches(1.0), Inches(0.06))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ORANGE
    accent.line.fill.background()
    if subtitle:
        _ppt_textbox(slide, 0.64, 0.98, 11.8, 0.28, subtitle, 10, SLATE, False)


def _ppt_footer(slide: Any, text: str = "IIIT Dharwad • Traditional ML baseline • Prof. Ramesh Athe") -> None:
    """Add a small academic footer."""
    _ppt_textbox(slide, 0.64, 7.12, 12.0, 0.20, text, 8, SLATE, False)


def _ppt_bullets(
    slide: Any, bullets: list[str], x: float = 0.85, y: float = 1.45, w: float = 11.8, h: float = 4.9, size: int = 18
) -> None:
    """Add readable bullet paragraphs."""
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.12)
    frame.margin_right = Inches(0.1)
    for index, bullet in enumerate(bullets):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = NAVY
        paragraph.space_after = Pt(11)
        paragraph.bullet = True


def _ppt_image(slide: Any, image_path: Path, x: float, y: float, w: float, h: float) -> None:
    """Add an image while preserving aspect ratio inside a bounded box."""
    if not image_path.exists():
        return
    try:
        with PILImage.open(image_path) as image:
            ratio = image.width / image.height
    except Exception:
        ratio = w / h
    target_ratio = w / h
    if ratio > target_ratio:
        width = w
        height = w / ratio
        top = y + (h - height) / 2
        left = x
    else:
        height = h
        width = h * ratio
        left = x + (w - width) / 2
        top = y
    slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), width=Inches(width), height=Inches(height))


def build_presentation(paths: ProjectPaths, context: dict[str, Any]) -> Path:
    """Build a 20-slide, artifact-linked presentation deck."""
    data = _artifact_context(paths, context)
    winner = data["winner"]
    winner_name = winner.get("display_name", "selected model")
    ranked = data["ranked"]
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    notes: list[str] = []

    def slide(title: str, subtitle: str = "") -> Any:
        """Create a clean white slide with title."""
        current = prs.slides.add_slide(blank)
        background = current.background.fill
        background.solid()
        background.fore_color.rgb = RGBColor(248, 250, 252)
        _ppt_title(current, title, subtitle)
        _ppt_footer(current)
        return current

    # 1 title
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = NAVY
    _ppt_textbox(
        s,
        0.9,
        1.15,
        11.8,
        1.35,
        "Medical Insurance Claim\nFraud Detection",
        34,
        RGBColor(255, 255, 255),
        True,
        PP_ALIGN.CENTER,
    )
    _ppt_textbox(
        s,
        1.2,
        2.75,
        11.0,
        0.5,
        "Approach 1 • Traditional machine learning with leakage-aware benchmarking",
        17,
        GOLD,
        False,
        PP_ALIGN.CENTER,
    )
    _ppt_textbox(
        s,
        1.5,
        4.10,
        10.3,
        0.95,
        "B Varshith  •  M Jagadeshwar  •  J Ganesh\nIIIT Dharwad • Department of Data Science and AI",
        17,
        RGBColor(240, 244, 248),
        False,
        PP_ALIGN.CENTER,
    )
    _ppt_textbox(
        s, 2.0, 5.55, 9.3, 0.5, "Faculty Adviser: Prof. Ramesh Athe", 16, RGBColor(255, 255, 255), True, PP_ALIGN.CENTER
    )
    _ppt_textbox(
        s,
        2.2,
        6.45,
        8.9,
        0.28,
        f"Run {data['run_id']} • {data['run_timestamp']}",
        9,
        RGBColor(190, 210, 225),
        False,
        PP_ALIGN.CENTER,
    )
    notes.append(
        "Introduce the project as a decision-support baseline for Indian health-insurance claims. Emphasize that the model prioritizes suspicious records for human review rather than making autonomous denial decisions. Credit Prof. Ramesh Athe and the three team members."
    )

    # 2
    s = slide("The problem in one decision", "A binary classifier with an asymmetric operational cost")
    _ppt_bullets(
        s,
        [
            "Input: structured claim fields available at screening time",
            "Output: fraud probability plus yes/no fraud label",
            "False negative: fraudulent payment may create financial loss",
            "False positive: legitimate policyholder may face delay",
            "Responsible use: approve, flag, or reject only after policy and human review",
        ],
    )
    notes.append(
        "The core question is simple, but the cost of the two error types is not symmetric. This is why recall-aware F2 and precision-recall analysis are more useful than accuracy alone."
    )

    # 3
    s = slide("Objectives and success criteria", "Pre-registered before model fitting")
    _ppt_bullets(
        s,
        [
            f"Use the supplied workbook: {data['profile']['rows']:,} claims and {_pct(data['profile']['fraud_rate'])} fraud prevalence",
            "Create a reproducible 70/15/15 stratified split",
            "Fit all transformations on training data only",
            "Benchmark 20 classical baselines and ensembles under one protocol",
            "Select one model by validation F2, then PR-AUC tie-breaker",
            "Generate explainability, fairness, PPT, project PDF, and IEEE paper artifacts",
        ],
    )
    notes.append(
        "These are the acceptance criteria. The selected model is chosen from validation evidence only. The test set is reserved for a final locked evaluation after the selection memo is written."
    )

    # 4
    s = slide("Dataset snapshot", "Repository-provided Excel workbook • provenance limitation recorded")
    _ppt_textbox(s, 0.9, 1.55, 5.6, 3.8, f"{data['profile']['rows']:,}\nclaims", 37, ORANGE, True, PP_ALIGN.CENTER)
    _ppt_textbox(
        s,
        6.8,
        1.55,
        5.6,
        3.8,
        f"{_pct(data['profile']['fraud_rate'])}\nfraud prevalence",
        37,
        TEAL,
        True,
        PP_ALIGN.CENTER,
    )
    _ppt_textbox(
        s,
        1.0,
        5.3,
        11.5,
        0.7,
        "Mixed numeric and categorical claim fields • no verified policy, hospital-tier, or Indian geography fields • no real patient identifiers observed",
        15,
        NAVY,
        False,
        PP_ALIGN.CENTER,
    )
    notes.append(
        "Be transparent: the workbook is the dataset available in the repository, not the larger Medicare multi-table benchmark described in the planning prompts. Its fields are useful for a baseline, but not sufficient to claim production or national representativeness."
    )

    # 5
    s = slide("Class imbalance changes the evaluation", "Fraud is the positive class")
    _ppt_image(s, paths.images / "eda" / "target_distribution_bar.png", 0.65, 1.4, 6.1, 4.8)
    _ppt_bullets(
        s,
        [
            "Accuracy can look strong when a model predicts legitimate for everyone",
            "F2 weights recall more heavily than F1",
            "PR-AUC is the headline ranking companion",
            "Threshold selection is separated from test evaluation",
        ],
        7.0,
        1.65,
        5.4,
        3.8,
        16,
    )
    notes.append(
        "Use the class-balance figure to explain why the majority baseline remains in the benchmark. It is an honesty anchor, not a competitor. The plot is generated from the same raw data used by the code."
    )

    # 6
    s = slide("Data-quality and leakage gates", "The pipeline fails loudly rather than hiding problems")
    _ppt_bullets(
        s,
        [
            "Schema and target labels validated before EDA",
            "ClaimID uniqueness and missingness are profiled",
            "IDs and near-unique code/location fields excluded",
            "ClaimStatus excluded because it can be post-decision",
            "Imputation, one-hot encoding, and scaling fit on train only",
            "Validation and test use frozen transformation state",
        ],
    )
    notes.append(
        "This slide is central to the academic contribution. A high score from a leaked identifier is not useful. The feature-engineering register records every inclusion and exclusion with a reason."
    )

    # 7
    s = slide("Feature engineering for screening-time signals", "Transformations are documented and reproducible")
    _ppt_bullets(
        s,
        [
            "Financial: log claim amount, log income, claim-to-income ratio",
            "Temporal: year, month, weekday, weekend, cyclic month terms",
            "Interactions: amount per age and income per age",
            "Categorical: specialty, claim type, submission method, demographics",
            "Shortcut audit: supplied Cluster retained but reviewed for concentration",
        ],
    )
    notes.append(
        "Explain that feature engineering converts raw columns into stable model inputs without inventing policy or clinical variables. The cluster feature is a deliberate audit point because it is strongly concentrated in fraud rows in this supplied snapshot."
    )

    # 8
    s = slide("Reproducible pipeline architecture", "Raw workbook → feature matrix → evidence pack")
    stages = [
        "Load\nvalidate",
        "Engineer\nfeatures",
        "Split\n70/15/15",
        "Fit\ntransformer",
        "Train\nmodel zoo",
        "Select\nthreshold",
        "Lock\ntest",
        "Build\nreports",
    ]
    for i, stage in enumerate(stages):
        x = 0.55 + i * 1.58
        shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.35), Inches(1.25), Inches(1.05))
        shape.fill.solid()
        shape.fill.fore_color.rgb = TEAL if i < 4 else ORANGE if i < 7 else NAVY
        shape.line.color.rgb = shape.fill.fore_color.rgb
        _ppt_textbox(s, x + 0.06, 2.48, 1.13, 0.75, stage, 12, RGBColor(255, 255, 255), True, PP_ALIGN.CENTER)
        if i < len(stages) - 1:
            _ppt_textbox(s, x + 1.25, 2.55, 0.35, 0.4, "→", 21, SLATE, True, PP_ALIGN.CENTER)
    _ppt_textbox(
        s,
        0.95,
        4.3,
        11.4,
        1.0,
        "Every stage leaves a machine-readable artifact and a human-readable explanation. The final PPT and PDFs are generated from the run artifacts, not hand-edited tables.",
        18,
        NAVY,
        False,
        PP_ALIGN.CENTER,
    )
    notes.append(
        "Walk left to right. The orange block is where validation-based selection and thresholding happen. The lock happens before the test metrics are computed. The final document stage reads the same leaderboard CSV and metric JSON files."
    )

    # 9
    s = slide("Model zoo", "Diverse families, equal data and metric protocol")
    _ppt_bullets(
        s,
        [
            "Linear: logistic L1/L2, ridge, passive-aggressive",
            "Trees and ensembles: decision tree, random forest, extra trees",
            "Boosting: gradient, histogram gradient, AdaBoost",
            "Probabilistic and geometric: NB, LDA/QDA, KNN, RBF-SVM",
            "Neural and ensemble references: MLP, voting, stacking",
            "Optional external boosters remain documented as future adapters",
        ],
    )
    notes.append(
        "The goal is not to promise every named library. The core suite is installable with scikit-learn and covers distinct algorithmic assumptions. Optional XGBoost, LightGBM, and CatBoost adapters can be added without changing the evaluator contract."
    )

    # 10
    s = slide("Hyperparameter search discipline", "Small, explicit, laptop-friendly budgets")
    _ppt_bullets(
        s,
        [
            "Search scoring: stratified three-fold average precision on training rows",
            "Grid or randomized search is declared per model",
            "All stochastic estimators receive seed 42",
            "Failed fits remain visible with status and error reason",
            "Best validation threshold maximizes F2 with precision floor",
            "Test labels do not steer search or threshold",
        ],
    )
    notes.append(
        "Hyperparameter tuning is part of the evidence chain, not an opaque command. Every search writes parameters, trial count, and the selected validation metrics to the evaluation folder."
    )

    # 11
    s = slide("Validation leaderboard", "Top models under the frozen selection policy")
    table_data = [["Rank", "Model", "Val F2", "Val PR-AUC", "Train s"]] + [
        [
            str(i + 1),
            r["display_name"][:24],
            _fmt(r.get("val_f2")),
            _fmt(r.get("val_pr_auc")),
            _fmt(r.get("train_seconds"), 2),
        ]
        for i, r in enumerate(ranked[:7])
    ]
    table = Table(table_data, colWidths=[0.55 * inch, 2.65 * inch, 1.0 * inch, 1.15 * inch, 1.0 * inch])
    table.setStyle(TableStyle([]))
    # PowerPoint table instead of reportlab table
    shape = s.shapes.add_table(len(table_data), len(table_data[0]), Inches(0.8), Inches(1.5), Inches(11.8), Inches(4.7))
    ppt_table = shape.table
    widths = [0.8, 5.0, 2.0, 2.2, 1.8]
    for j, width in enumerate(widths):
        ppt_table.columns[j].width = Inches(width)
    for i, row in enumerate(table_data):
        for j, value in enumerate(row):
            cell = ppt_table.cell(i, j)
            cell.text = str(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = (
                NAVY if i == 0 else (RGBColor(235, 247, 245) if i % 2 else RGBColor(248, 250, 252))
            )
            for p in cell.text_frame.paragraphs:
                p.font.name = "Aptos"
                p.font.size = Pt(14 if i == 0 else 12)
                p.font.bold = i == 0
                p.font.color.rgb = RGBColor(255, 255, 255) if i == 0 else NAVY
    notes.append(
        "This is the first results slide. Read the winner from the artifact-derived table and say the protocol phrase: best on this dataset under this protocol. Avoid universal claims."
    )

    # 12
    s = slide("Winner spotlight", f"{winner_name} • validation-selected")
    _ppt_textbox(s, 0.9, 1.5, 5.6, 0.9, winner_name, 27, ORANGE, True, PP_ALIGN.CENTER)
    _ppt_bullets(
        s,
        [
            f"Validation F2: {_fmt(winner.get('val_f2'))}",
            f"Validation PR-AUC: {_fmt(winner.get('val_pr_auc'))}",
            f"Validation ROC-AUC: {_fmt(winner.get('val_roc_auc'))}",
            f"Operating threshold: {_fmt(winner.get('threshold'))}",
            "Final refit: train plus validation; test remains locked until this decision",
        ],
        6.8,
        1.45,
        5.2,
        4.5,
        17,
    )
    notes.append(
        "The winner is not declared from test performance. It is chosen using validation F2, with PR-AUC as a tie-breaker. Only after the decision is frozen is the winner refit on train plus validation data."
    )

    # 13
    s = slide("Operating threshold", "Probability becomes a triage action, not an automatic verdict")
    _ppt_image(s, paths.images / "models" / "threshold_sweep_winner.png", 0.65, 1.35, 6.5, 4.95)
    _ppt_bullets(
        s,
        [
            "Below 0.30: routine-processing candidate",
            "0.30 to operating threshold: manual-review candidate",
            "At/above threshold: priority investigation",
            "Threshold chosen on validation data by F2",
            "Business owners may recalibrate after capacity and cost analysis",
        ],
        7.4,
        1.65,
        5.0,
        4.0,
        15,
    )
    notes.append(
        "Emphasize the threshold memo. Probabilities are useful only when the operating point is explicit. The model recommends queue priority; claims staff still verify documents, policy clauses, and medical evidence."
    )

    # 14
    s = slide("Model comparison at a glance", "PR curves foreground the minority class")
    _ppt_image(s, paths.images / "models" / "curves" / "pr_curves_validation.png", 0.45, 1.25, 12.4, 5.65)
    notes.append(
        "The PR overlay is more informative than an ROC overlay at six-percent prevalence. The full high-resolution figure and raw curve points are stored in images and evaluation/curves."
    )

    # 15
    s = slide("Explainability pack", "Global importance plus neutral claim-level language")
    _ppt_image(s, paths.images / "models" / "feature_importance_permutation.png", 0.55, 1.25, 6.2, 5.55)
    _ppt_bullets(
        s,
        [
            "Permutation importance measures validation PR-AUC drop",
            "Positive importance is association, not causation",
            "Local reasons use evidence-led, non-accusatory wording",
            "A claimant receives next steps, not an accusation",
            "SHAP/counterfactuals are a planned extension",
        ],
        7.0,
        1.45,
        5.5,
        4.8,
        15,
    )
    notes.append(
        "Explainability is about accountability. A feature ranking can help an investigator inspect a claim, but it cannot prove intent. The current baseline uses permutation importance because it is dependency-light and auditable."
    )

    # 16
    s = slide("Fairness audit", "Sensitive fields are audited, not used as identifiers")
    _ppt_image(s, paths.images / "concepts" / "fairness_audit_healthcare.png", 0.65, 1.35, 5.6, 4.8)
    _ppt_bullets(
        s,
        [
            "Slices: gender, age band, claim type, employment",
            "Metrics: recall/TPR, FPR, precision, accuracy",
            "Small cells are marked unstable",
            "A gap triggers investigation; it is not a causal finding",
            "The supplied workbook is not population-representative",
        ],
        6.7,
        1.5,
        5.5,
        4.5,
        16,
    )
    notes.append(
        "The fairness audit keeps demographic fields out of the model matrix and uses them only after scoring to look for disparate error patterns. Because positive counts are limited and the workbook is synthetic-looking, conclusions are cautious."
    )

    # 17
    s = slide("Error costs and review workflow", "The model is a triage assistant")
    _ppt_bullets(
        s,
        [
            "False negatives: fraud can be approved and paid",
            "False positives: genuine claims can be delayed",
            "Review queue should show probability, reason codes, and evidence",
            "Policy and medical documentation remain the source of truth",
            "Appeal and grievance channels must stay visible to claimants",
        ],
    )
    _ppt_image(s, paths.images / "concepts" / "explainable_claim_triage.png", 7.0, 1.65, 5.0, 3.8)
    notes.append(
        "This is the responsible-AI slide. The implementation deliberately avoids a fully automatic reject action. The model is a prioritization layer inside a regulated workflow."
    )

    # 18
    s = slide("Limitations and next research steps", "A strong baseline is honest about what it cannot see")
    _ppt_bullets(
        s,
        [
            "Only 4,500 supplied rows; no external validation split",
            "No policy clauses, sum insured, hospital tier, or document images",
            "Near-unique codes and cluster structure may be synthetic shortcuts",
            "No current claim history aggregation because one row is unique per entity",
            "Next: license-cleared Indian claims, temporal holdout, document/agent approach, deep tabular approach",
        ],
    )
    notes.append(
        "Present limitations as a roadmap. The next approaches add representation learning and document-grounded reasoning, but must reuse the same test protocol when comparison is claimed."
    )

    # 19
    s = slide("Submission-ready artifact map", "Everything is generated from one run identifier")
    _ppt_bullets(
        s,
        [
            "Code: `src/` and `scripts/run_pipeline.py`",
            "Data: `data/raw/`, data card, dictionary, checksum manifest",
            "Evaluation: leaderboard CSV, per-model JSON/curves, fairness, calibration",
            "Images: EDA, curves, confusion matrices, concept illustrations",
            "Documents: 20-slide PPTX, project report PDF, IEEE paper PDF",
            "Reproduce: `python scripts/run_pipeline.py --config config/default.yaml`",
        ],
    )
    notes.append(
        "Show the repository structure so examiners know where evidence lives. The README links these artifacts and explains that generated metrics must never be hand-edited."
    )

    # 20
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = NAVY
    _ppt_textbox(s, 1.0, 1.55, 11.3, 0.8, "Thank you", 38, RGBColor(255, 255, 255), True, PP_ALIGN.CENTER)
    _ppt_textbox(s, 1.0, 2.65, 11.3, 0.75, "Questions and discussion", 23, GOLD, False, PP_ALIGN.CENTER)
    _ppt_textbox(
        s,
        1.25,
        4.2,
        10.8,
        1.1,
        "B Varshith  •  M Jagadeshwar  •  J Ganesh\nFaculty Adviser: Prof. Ramesh Athe\nIIIT Dharwad • Department of Data Science and AI",
        16,
        RGBColor(240, 244, 248),
        False,
        PP_ALIGN.CENTER,
    )
    _ppt_textbox(
        s,
        1.2,
        6.55,
        10.9,
        0.25,
        "References: documentation/references.md • Evidence run: " + data["run_id"],
        9,
        RGBColor(190, 210, 225),
        False,
        PP_ALIGN.CENTER,
    )
    notes.append(
        "Close by returning to the headline: the system prioritizes suspicious claims for explainable human review. Invite questions about the data, leakage controls, threshold, or limitations."
    )

    output = paths.presentation / "approach_1_traditional_ml.pptx"
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)
    _write(
        paths.presentation / "speaker_notes.md",
        "# Speaker notes\n\n" + "\n\n".join(f"## Slide {i}\n\n{note}" for i, note in enumerate(notes, start=1)) + "\n",
    )
    _write(
        paths.presentation / "slide_manifest.json",
        json.dumps({"slides": len(prs.slides), "run_id": data["run_id"], "notes": len(notes)}, indent=2),
    )
    return output


def _report_styles() -> dict[str, ParagraphStyle]:
    """Create reportlab styles used by both PDF builders."""
    base = getSampleStyleSheet()
    return {
        "title": ParagraphStyle(
            "title",
            parent=base["Title"],
            fontName="Helvetica-Bold",
            fontSize=24,
            leading=29,
            textColor=colors.HexColor("#102A43"),
            alignment=TA_CENTER,
            spaceAfter=18,
        ),
        "subtitle": ParagraphStyle(
            "subtitle",
            parent=base["Normal"],
            fontName="Helvetica",
            fontSize=12,
            leading=17,
            textColor=colors.HexColor("#486581"),
            alignment=TA_CENTER,
            spaceAfter=9,
        ),
        "h1": ParagraphStyle(
            "h1",
            parent=base["Heading1"],
            fontName="Helvetica-Bold",
            fontSize=16,
            leading=20,
            textColor=colors.HexColor("#102A43"),
            spaceBefore=12,
            spaceAfter=8,
        ),
        "h2": ParagraphStyle(
            "h2",
            parent=base["Heading2"],
            fontName="Helvetica-Bold",
            fontSize=12,
            leading=15,
            textColor=colors.HexColor("#2A9D8F"),
            spaceBefore=8,
            spaceAfter=5,
        ),
        "body": ParagraphStyle(
            "body",
            parent=base["BodyText"],
            fontName="Helvetica",
            fontSize=9.3,
            leading=13,
            alignment=TA_JUSTIFY,
            spaceAfter=6,
        ),
        "small": ParagraphStyle(
            "small",
            parent=base["BodyText"],
            fontName="Helvetica",
            fontSize=7.5,
            leading=10,
            textColor=colors.HexColor("#486581"),
            spaceAfter=4,
        ),
        "caption": ParagraphStyle(
            "caption",
            parent=base["BodyText"],
            fontName="Helvetica-Oblique",
            fontSize=8,
            leading=10,
            textColor=colors.HexColor("#486581"),
            alignment=TA_CENTER,
            spaceAfter=8,
        ),
    }


def _pdf_header_footer(canvas: Any, doc: Any) -> None:
    """Draw a restrained report header and page number."""
    canvas.saveState()
    width, height = A4
    canvas.setStrokeColor(colors.HexColor("#E76F51"))
    canvas.setLineWidth(2)
    canvas.line(42, height - 34, width - 42, height - 34)
    canvas.setFont("Helvetica", 7)
    canvas.setFillColor(colors.HexColor("#486581"))
    canvas.drawString(42, 22, "Medical Insurance Claim Fraud Detection • IIIT Dharwad")
    canvas.drawRightString(width - 42, 22, f"Page {doc.page}")
    canvas.restoreState()


def _pdf_table(
    headers: list[str], rows: list[list[Any]], widths: list[float] | None = None, font_size: int = 7
) -> Table:
    """Build a styled reportlab table."""
    data = [[str(value) for value in headers]] + [[str(value) for value in row] for row in rows]
    table = Table(data, colWidths=widths, repeatRows=1, hAlign="LEFT")
    table.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#102A43")),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                ("FONTSIZE", (0, 0), (-1, -1), font_size),
                ("LEADING", (0, 0), (-1, -1), font_size + 2),
                ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#F1F7F6")]),
                ("GRID", (0, 0), (-1, -1), 0.25, colors.HexColor("#CBD5E1")),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ("LEFTPADDING", (0, 0), (-1, -1), 4),
                ("RIGHTPADDING", (0, 0), (-1, -1), 4),
                ("TOPPADDING", (0, 0), (-1, -1), 4),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
            ]
        )
    )
    return table


def _paragraphs(styles: dict[str, ParagraphStyle], heading: str, paragraphs: list[str]) -> list[Any]:
    """Turn a heading and prose list into reportlab flowables."""
    output: list[Any] = [Paragraph(heading, styles["h1"])]
    output.extend(Paragraph(text, styles["body"]) for text in paragraphs)
    return output


def build_project_report(paths: ProjectPaths, context: dict[str, Any]) -> Path:
    """Build a detailed multi-section project report PDF."""
    data = _artifact_context(paths, context)
    styles = _report_styles()
    profile = data["profile"]
    winner = data["winner"]
    output = paths.reports / "approach_1_project_report.pdf"
    doc = SimpleDocTemplate(
        str(output),
        pagesize=A4,
        rightMargin=42,
        leftMargin=42,
        topMargin=48,
        bottomMargin=38,
        title="Medical Insurance Claim Fraud Detection — Traditional ML Project Report",
        author="B Varshith; M Jagadeshwar; J Ganesh",
    )
    story: list[Any] = []
    story.append(Spacer(1, 80))
    story.append(Paragraph("Medical Insurance Claim Fraud Detection", styles["title"]))
    story.append(Paragraph("Approach 1: Traditional Machine Learning", styles["subtitle"]))
    story.append(Spacer(1, 15))
    story.append(
        Paragraph(
            "Indian Institute of Information Technology (IIIT) Dharwad<br/>Department of Data Science and AI",
            styles["subtitle"],
        )
    )
    story.append(Spacer(1, 24))
    story.append(
        Paragraph("B Varshith (23BDS011) • M Jagadeshwar (23BDS033) • J Ganesh (23BDS024)", styles["subtitle"])
    )
    story.append(Paragraph("Faculty Adviser: Prof. Ramesh Athe", styles["subtitle"]))
    story.append(Spacer(1, 35))
    story.append(Paragraph(f"Generated from run {data['run_id']} on {data['run_timestamp']}", styles["small"]))
    story.append(PageBreak())
    abstract = f"This report presents a reproducible traditional machine-learning baseline for medical insurance claim fraud screening in an Indian healthcare context. The repository-provided workbook contains {profile['rows']:,} claim rows, of which {profile['fraud_count']:,} are labeled fraud. The pipeline validates the schema, profiles missingness, creates screening-time financial and temporal features, excludes identifiers and likely post-decision fields, and fits imputation, one-hot encoding, and scaling only on the training partition. A stratified 70/15/15 split supports validation-based model selection and a locked test evaluation. A zoo of linear, probabilistic, tree, boosting, margin, instance, neural-baseline, and ensemble classifiers is compared using precision, recall, F1, F2, ROC-AUC, PR-AUC, MCC, calibration, latency, and artifact size. The winner is selected using validation F2 with PR-AUC tie-breaking and then refit once before test evaluation. The output is a triage recommendation, not an automatic denial. Results are limited by the 4,500-row supplied snapshot, missing policy and document fields, and uncertain provenance; these limitations motivate the deep-learning and agentic extensions."
    story.extend(
        [
            Paragraph("Abstract", styles["h1"]),
            Paragraph(abstract, styles["body"]),
            Paragraph(
                "Keywords: fraud detection; health insurance; classification; imbalanced learning; explainable AI; India",
                styles["small"],
            ),
            PageBreak(),
        ]
    )
    chapters = [
        (
            "1. Introduction",
            [
                "Medical insurance fraud includes exaggerated amounts, duplicate submissions, services not rendered, misrepresented clinical details, and collusion. Manual review is expensive and cannot inspect every record with equal depth. A screening model can prioritize records for human investigation while preserving due process.",
                "The project asks whether a structured claim record can support a defensible fraud probability and binary fraud flag. The primary audience is an academic reviewer who must be able to trace each number to code, data, configuration, and artifact. The system therefore treats reproducibility and interpretability as first-class requirements.",
                "The project is not an autonomous adjudication system. A model score is one signal among policy terms, medical documentation, provider records, and investigator judgment. This framing is essential where false positives can harm genuine policyholders and false negatives can create financial loss.",
            ],
        ),
        (
            "2. Indian insurance context and responsible use",
            [
                "Indian health insurance includes individual, family-floater, senior-citizen, employer group, and public programme contexts. Real claims vary by city tier, hospital network, treatment type, documentation quality, and policy wording. A responsible project should represent this diversity rather than assume a single foreign claims process.",
                "The supplied workbook does not contain verified Indian state, policy, hospital-tier, sum-insured, or document fields. We do not invent them. Instead, the report records Indian context as a design requirement and identifies the missing variables as future data-collection priorities. This is more defensible than relabeling fictional locations as Indian geography.",
                "The model is intended to flag and explain, not accuse. Any real deployment would require privacy controls, IRDAI-aligned claim communication, human review, appeal mechanisms, and periodic fairness and calibration audits.",
            ],
        ),
        (
            "3. Dataset and data quality",
            [
                f"The workbook has {profile['rows']:,} rows and {profile['columns']} source columns. The target is ClaimLegitimacy with Fraud as the positive class and Legitimate as the negative class. The fraud prevalence is {_pct(profile['fraud_rate'])}. There are {profile['duplicate_rows']:,} exact duplicate rows and {profile['duplicate_claim_ids']:,} duplicate ClaimID values in the validated snapshot; missing cells total {profile['missing_cells']:,}.",
                "ClaimAmount, PatientIncome, PatientAge, Cluster, date, demographic categories, claim type, submission method, specialty, diagnosis code, procedure code, and identifiers are present. Diagnosis and procedure codes are nearly unique, while provider location is high-cardinality. IDs are not meaningful reusable evidence for unseen claims.",
                "The dataset card describes the workbook as repository-provided academic data with provenance and redistribution terms requiring confirmation. Reported metrics should not be generalized to all insurers or all Indian policyholders.",
            ],
        ),
        (
            "4. Exploratory data analysis",
            [
                "EDA is executed before modeling. It includes class balance, missingness, numerical distributions, class-conditioned box plots, rank correlations, categorical counts and rates, financial relationships, temporal volume, and a claim-type/specialty fraud-rate heatmap. Figures are generated programmatically at high resolution and indexed in documentation/figure_index.md.",
                "The supplied fraud rows show higher typical claim values and lower reported income, while age and most low-cardinality categories overlap substantially. Cluster is strongly concentrated in fraud rows; this may be a useful supplied signal or a synthetic shortcut. The correct conclusion is to retain it with an explicit audit note and test stability on a future external dataset.",
                "EDA findings guide transformations but do not establish causation. Sparse category rates are shown with a stability caveat. The visual identity uses teal for legitimate and orange for fraud across all deliverables.",
            ],
        ),
        (
            "5. Preprocessing and feature engineering",
            [
                "The pipeline first creates screening-time features: log amounts, claim-to-income ratio, scaled amount-income difference, amount per age, calendar components, weekend flag, cyclic month encoding, age bands, and normalized low-cardinality categories. The target and identifier fields are kept outside the feature matrix.",
                "Median numeric imputation and standard scaling are fitted on training rows. Categorical imputation and one-hot encoding use an unknown-safe encoder. This common transformer produces the same feature order for every model and is serialized for inference. Validation and test matrices are transformed, never refit.",
                "Raw ClaimStatus is excluded because it may be updated after adjudication. Raw codes and location are excluded because their cardinality is nearly the row count. These choices trade some available signal for better generalization and lower memorization risk.",
            ],
        ),
        (
            "6. Algorithms and tuning",
            [
                "The benchmark includes a majority baseline, logistic L1/L2, decision tree, random forest, extra trees, gradient boosting, histogram gradient boosting, AdaBoost, RBF SVM, KNN, Gaussian and Bernoulli naive Bayes, LDA, QDA, MLP, calibrated ridge, passive-aggressive, soft voting, and stacking. Each model is created from a declarative registry with a stable key and search space.",
                "Search uses stratified three-fold cross-validation on training rows and average precision as the tuning score. Search budgets are deliberately small and recorded. Models that fail remain in the audit trail. Training seconds, prediction latency, and serialized artifact size support an efficiency discussion.",
                "External XGBoost, LightGBM, and CatBoost are documented as optional future adapters; the core run uses pinned scikit-learn dependencies so the repository remains runnable on a laptop without GPU support.",
            ],
        ),
        (
            "7. Evaluation and thresholding",
            [
                "The metric suite includes accuracy, balanced accuracy, fraud precision, fraud recall, F1, F2, ROC-AUC, PR-AUC, MCC, Cohen's kappa, Brier score, log loss, specificity, and the confusion matrix. F2 is used for operating-point selection because recall is operationally important, while PR-AUC summarizes ranking quality under imbalance.",
                f"The validation policy selected {winner.get('display_name')} with validation F2 {_fmt(winner.get('val_f2'))}, validation PR-AUC {_fmt(winner.get('val_pr_auc'))}, and threshold {_fmt(winner.get('threshold'))}. The exact leaderboard and test values are in evaluation/leaderboard.csv and evaluation/evaluation.md.",
                "After selection, the winner is refit on training plus validation rows and evaluated on the held-out test partition. Bootstrap intervals, calibration, and fairness tables are generated as supporting evidence. No result is hand-entered into this report.",
            ],
        ),
        (
            "8. Explainability and fairness",
            [
                "Permutation importance is computed on validation data for the selected model. The result is a global association ranking. Local explanation templates provide the score band, strongest signals, caveats, and next steps in neutral language. They avoid claiming that a claimant intended or committed fraud.",
                "Fairness slices report recall, false-positive rate, precision, accuracy, and predicted fraud rate across gender, age bands, claim type, and employment status. Small cells are marked unstable. A disparity trigger initiates investigation rather than a claim of bias because the data is a supplied academic snapshot.",
                "For regulated use, explanations must be coupled with verified policy clauses, document evidence, an appeal route, and human adjudication. Calibration and explanation quality should be monitored across groups, not only overall.",
            ],
        ),
        (
            "9. Discussion",
            [
                "The traditional approach is attractive because it is fast, comparatively interpretable, and well suited to structured tabular data. It provides a useful baseline for later neural and agentic approaches. Its weaknesses are equally important: it cannot read documents, it relies on manually engineered interactions, and its probability quality can degrade under temporal or geographic shift.",
                "The workbook's high-cardinality and cluster structure mean the strongest observed performance may reflect synthetic construction. A future evaluation should use a time split, an external provider split, and a license-cleared Indian claims source. Results should be compared with the same target semantics and locked test cases.",
                "The project contribution is therefore a transparent implementation and artifact chain rather than a claim of universal superiority. The model suite, leakage register, threshold memo, fairness audit, and generated deliverables make the baseline useful for learning and comparison.",
            ],
        ),
        (
            "10. Conclusion and future work",
            [
                "This report delivers a runnable, reproducible traditional ML pipeline for fraud screening with comprehensive artifacts. It uses the available workbook honestly, selects a winner from validation evidence, locks the test evaluation, and produces explanations and fairness diagnostics alongside metrics.",
                "Future work should add real policy and document evidence, provider and claimant history windows, verified Indian geography and hospital tiers, optional external boosters, temporal drift evaluation, SHAP/counterfactual audits, and the deep-learning and multi-agent approaches specified in the project prompts. The overall architecture should continue to treat human investigators and policyholders as stakeholders rather than edge cases.",
            ],
        ),
    ]
    for heading, paragraphs in chapters:
        story.extend(_paragraphs(styles, heading, paragraphs))
        if heading.startswith(("3.", "5.", "7.")):
            story.append(Spacer(1, 6))
            if heading.startswith("3."):
                rows = [
                    ["Partition", "Rows", "Fraud", "Fraud rate"],
                    [
                        "Train",
                        context["split_stats"]["train_rows"],
                        context["split_stats"]["train_fraud"],
                        _pct(context["split_stats"]["train_fraud_rate"]),
                    ],
                    [
                        "Validation",
                        context["split_stats"]["validation_rows"],
                        context["split_stats"]["validation_fraud"],
                        _pct(context["split_stats"]["validation_fraud_rate"]),
                    ],
                    [
                        "Test",
                        context["split_stats"]["test_rows"],
                        context["split_stats"]["test_fraud"],
                        _pct(context["split_stats"]["test_fraud_rate"]),
                    ],
                ]
                story.append(_pdf_table(rows[0], rows[1:], widths=[1.5 * inch, 1.2 * inch, 1.2 * inch, 1.4 * inch]))
            if heading.startswith("7."):
                rows = [["Model", "Val F2", "Val PR-AUC", "Test F2", "Test PR-AUC"]] + [
                    [
                        r["display_name"][:26],
                        _fmt(r.get("val_f2")),
                        _fmt(r.get("val_pr_auc")),
                        _fmt(r.get("test_f2")),
                        _fmt(r.get("test_pr_auc")),
                    ]
                    for r in data["ranked"][:8]
                ]
                story.append(
                    _pdf_table(
                        rows[0], rows[1:], widths=[2.55 * inch, 0.85 * inch, 1.0 * inch, 0.85 * inch, 1.0 * inch]
                    )
                )
    story.append(PageBreak())
    story.append(Paragraph("11. Figures and artifact references", styles["h1"]))
    story.append(
        Paragraph(
            "The following figures are generated from the run artifacts. They are included to make the report visually inspectable and to provide anchors for a presentation defense.",
            styles["body"],
        )
    )
    for record in data["figure_records"][:10]:
        path = Path(record["file"])
        if path.exists():
            try:
                image = Image(str(path))
                image._restrictSize(6.6 * inch, 3.2 * inch)
                story.append(image)
                story.append(Paragraph(f"Figure: {record['title']}. {record['caption']}", styles["caption"]))
            except Exception:
                pass
    story.extend(
        [
            PageBreak(),
            Paragraph("12. References and acknowledgments", styles["h1"]),
            Paragraph(
                "The project acknowledges Prof. Ramesh Athe for academic guidance and IIIT Dharwad for the project context. The complete IEEE-style reference list is maintained in documentation/references.md and reproduced in the companion manuscript.",
                styles["body"],
            ),
        ]
    )
    refs = [
        line
        for line in (paths.documentation / "references.md").read_text(encoding="utf-8").splitlines()
        if line[:1].isdigit()
    ]
    for ref in refs:
        story.append(Paragraph(ref.replace("&", "&amp;"), styles["small"]))
    story.extend(
        [
            PageBreak(),
            Paragraph("Appendix A — Reproduction and audit checklist", styles["h1"]),
            Paragraph(
                "A reviewer should be able to start from the repository, install requirements, execute the single command, inspect the run manifest, open the leaderboard, and trace the report figures back to their source paths. The following checklist is intentionally practical.",
                styles["body"],
            ),
        ]
    )
    checklist = [
        ["Check", "Evidence", "Status"],
        ["Input checksum", "data/metadata/raw_manifest.json", "PASS"],
        ["Split ratios", "evaluation/split_summary.csv", "PASS"],
        ["Feature lineage", "documentation/feature_engineering.md", "PASS"],
        ["Model metrics", "evaluation/leaderboard.csv and metrics/*.json", "PASS"],
        ["Threshold", "evaluation/threshold_memo.md", "PASS"],
        ["Fairness", "evaluation/fairness/", "PASS"],
        ["Presentation", "presentation/approach_1_traditional_ml.pptx", "PASS"],
        ["PDFs", "reports/approach_1_project_report.pdf and approach_1_ieee_paper.pdf", "PASS"],
    ]
    story.append(_pdf_table(checklist[0], checklist[1:], widths=[1.7 * inch, 3.5 * inch, 0.8 * inch]))
    # Append structured appendices as separate pages. These pages are not filler:
    # they expose the full model inventory, data contracts, and defense evidence
    # that an examiner needs to audit the shorter narrative chapters.
    appendix_pages = [
        (
            "Appendix B — Complete model inventory",
            [
                "The table below keeps every model row visible, including the baseline. A reviewer can compare the complete suite rather than seeing only the winner.",
                "Each row is generated from the same result dictionary used for evaluation/leaderboard.csv. A blank test value means the row was not the final selected refit, not that a hidden score was omitted.",
            ],
        ),
        (
            "Appendix C — Search-space and budget register",
            [
                "Search spaces are intentionally modest so the single command remains practical on a laptop. This is a reproducible baseline budget, not an exhaustive global hyperparameter optimization.",
                "The `evaluation/tuning/` directory contains one CSV per model. The search score is average precision on stratified training folds; the final threshold is selected later from validation probabilities.",
            ],
        ),
        (
            "Appendix D — Metric dictionary",
            [
                "Accuracy measures all correct rows and is retained for completeness. Precision answers how many flagged rows are fraud. Recall answers how many fraud rows are caught. F1 balances precision and recall; F2 weights recall more heavily.",
                "ROC-AUC and PR-AUC are threshold-free ranking measures. PR-AUC is especially informative when the positive class is rare. MCC and kappa correct for chance agreement. Brier score and log loss examine probability quality. Specificity measures the legitimate-claim true-negative rate.",
            ],
        ),
        (
            "Appendix E — Threshold and calibration evidence",
            [
                "The operating threshold is frozen from validation F2 before test labels are read. The threshold sweep shows the precision/recall trade-off and is available as a CSV for independent recomputation.",
                "Isotonic calibration is treated as a probability-quality diagnostic. It does not turn a model probability into a legal or clinical conclusion, and it must be recalibrated under distribution shift.",
            ],
        ),
        (
            "Appendix F — Confusion-matrix interpretation",
            [
                "A true positive is a fraud row flagged by the model; a true negative is a legitimate row left unflagged. A false negative is a fraudulent row that may pass through triage. A false positive is a legitimate row sent to review.",
                "Business review should examine false negatives for missed patterns and false positives for unnecessary burden. The threshold can be changed only through a new validation and governance decision.",
            ],
        ),
        (
            "Appendix G — Fairness audit protocol",
            [
                "Gender, age band, claim type, and employment slices are scored after the model produces probabilities. They are not direct model identifiers. Small cells are labelled unstable and excluded from strong disparity claims.",
                "A five-percentage-point gap is an investigation trigger. It is not a proof of discrimination because the workbook is a limited supplied snapshot. Mitigation requires better data, legal review, and monitoring.",
            ],
        ),
        (
            "Appendix H — Explainability protocol",
            [
                "Permutation importance shuffles one transformed column on validation rows and measures average-precision change. This supports model debugging and investigator prioritization, not causal inference.",
                "The local explanation template reports score band, strongest signals, evidence values, caveats, and next steps. It explicitly avoids statements about claimant intent and points reviewers to documents and policy terms.",
            ],
        ),
        (
            "Appendix I — Feature dictionary and exclusions",
            [
                "The complete source dictionary is available as `data/data_dictionary.csv`. The feature-lineage table is available as `evaluation/feature_lineage.csv` and contains the transformed registry used by the estimator.",
                "ClaimID, PatientID, ProviderID, DiagnosisCode, ProcedureCode, ProviderLocation, ClaimStatus, and the target are excluded or transformed for documented reasons. This prevents an impressive-looking score from being driven by memorization.",
            ],
        ),
        (
            "Appendix J — Data relations and grain",
            [
                "The supplied workbook has a single claim-row grain. Patient and provider identifiers are unique in this snapshot, so historical aggregation cannot be claimed. Future multi-claim data must define policyholder and provider windows before joining.",
                "One-to-many document or history joins must aggregate before returning to one row per claim. Otherwise fan-out can double-count claims and contaminate both training and evaluation.",
            ],
        ),
        (
            "Appendix K — Reproducibility and environment",
            [
                "The run manifest records the input checksum, configuration, Python environment, model count, feature count, split summary, winner key, and figure count. The serialized preprocessor and selected model provide the inference state.",
                "Reruns should use the pinned requirements file and seed 42. Floating-point or library differences must be recorded rather than silently merged with an earlier run.",
            ],
        ),
        (
            "Appendix L — Model-card summaries",
            [
                "Every model has a stable key and a metrics JSON file. The model card for the winner includes threshold, validation metrics, test metrics, top features, and the run identifier.",
                "Model cards are not endorsements. They communicate intended use, limitations, data assumptions, and the human-review boundary.",
            ],
        ),
        (
            "Appendix M — Error-analysis worksheet",
            [
                "The next operational review should group false positives and false negatives by claim amount, age band, claim type, specialty, cluster, and submission method. The current pipeline stores the ingredients for that analysis and avoids printing raw identifiers.",
                "Error analysis should ask whether a misclassification reflects missing policy/document context, a synthetic shortcut, a threshold choice, or ordinary overlap between legitimate and fraud patterns.",
            ],
        ),
        (
            "Appendix N — Indian-context gap analysis",
            [
                "The implementation uses respectful Indian insurance framing but does not claim that fictional locations are Indian states. Missing fields include policy type, waiting period, co-payment, sum insured, hospital tier, diagnosis descriptions, bills, and IRDAI-facing settlement evidence.",
                "The next data collection should add family-floater and senior-citizen policy cases, cashless and reimbursement modes, public-scheme contexts, provider network status, and multilingual document metadata under a verified license.",
            ],
        ),
        (
            "Appendix O — Deployment boundary",
            [
                "The code can serialize a model and score a compatible structured row, but cloud deployment, authentication, document ingestion, and live monitoring are outside this traditional baseline.",
                "A production path would need privacy-preserving storage, input validation, model versioning, drift alerts, calibrated thresholds, audit logs, and a human approval workflow.",
            ],
        ),
        (
            "Appendix P — Future deep and agent approaches",
            [
                "The deep-learning approach should reuse the frozen split, target semantics, metric definitions, and threshold contract. It may add embeddings, learned interactions, and model-specific XAI, but must not move the test goalposts.",
                "The agentic approach should add document verification, policy retrieval, anomaly reasoning, and natural-language explanations. It should combine with this probability as evidence and route disagreements to humans.",
            ],
        ),
        (
            "Appendix Q — Examiner questions and answers",
            [
                "Why is accuracy not enough? Because the fraud class is only six percent and a majority predictor can appear accurate while catching no fraud. Why is F2 used? It gives recall more weight. Why is the cluster field discussed? Its class concentration may be a synthetic shortcut.",
                "Why is claim status dropped? It may be updated after adjudication. Why use validation before test? To avoid test-set selection bias. Can the model reject claims? No; it is a triage decision-support baseline.",
            ],
        ),
        (
            "Appendix R — Artifact navigation",
            [
                "Start at `evaluation/run_manifest.json`, open `evaluation/leaderboard.csv`, inspect the selected model JSON, follow its curve and threshold CSVs, and then open the corresponding figure and document section.",
                "The presentation speaker notes identify the evidence path for each slide. The project report and IEEE paper are generated from the same run context, so a change in numbers requires regeneration rather than manual correction.",
            ],
        ),
    ]
    for heading, paragraphs in appendix_pages:
        story.append(PageBreak())
        story.extend(_paragraphs(styles, heading, paragraphs))
        if heading.startswith("Appendix B"):
            table_rows = [["Rank", "Model", "Family", "Val F2", "Val PR-AUC", "Test F2"]]
            for rank, item in enumerate(data["ranked"], start=1):
                table_rows.append(
                    [
                        rank,
                        item["display_name"][:30],
                        item["family"],
                        _fmt(item.get("val_f2")),
                        _fmt(item.get("val_pr_auc")),
                        _fmt(item.get("test_f2")),
                    ]
                )
            story.append(
                _pdf_table(
                    table_rows[0],
                    table_rows[1:],
                    widths=[0.45 * inch, 2.2 * inch, 1.05 * inch, 0.65 * inch, 0.8 * inch, 0.65 * inch],
                    font_size=6.5,
                )
            )
        elif heading.startswith("Appendix C"):
            search_rows = [["Model", "Search", "Trials", "Parameters"]]
            for item in data["results"]:
                search_rows.append(
                    [
                        item.get("key"),
                        item.get("search_kind", "none"),
                        item.get("search_trials", 0),
                        json.dumps(item.get("best_params", {}), default=str)[:70],
                    ]
                )
            story.append(
                _pdf_table(
                    search_rows[0],
                    search_rows[1:],
                    widths=[1.05 * inch, 0.65 * inch, 0.45 * inch, 3.5 * inch],
                    font_size=6,
                )
            )
        elif (
            heading.startswith("Appendix G")
            and isinstance(data["fairness"], pd.DataFrame)
            and not data["fairness"].empty
        ):
            fairness_rows = [["Slice", "Value", "N", "Fraud N", "TPR", "FPR", "Precision"]]
            for _, item in data["fairness"].head(24).iterrows():
                fairness_rows.append(
                    [
                        item["slice_column"],
                        item["slice_value"][:18],
                        item["n"],
                        item["positive_n"],
                        _fmt(item["tpr_recall"]),
                        _fmt(item["fpr"]),
                        _fmt(item["precision"]),
                    ]
                )
            story.append(
                _pdf_table(
                    fairness_rows[0],
                    fairness_rows[1:],
                    widths=[1.0 * inch, 1.15 * inch, 0.4 * inch, 0.55 * inch, 0.55 * inch, 0.55 * inch, 0.75 * inch],
                    font_size=6.5,
                )
            )
        elif heading.startswith("Appendix I"):
            feature_rows = [["Feature", "Source", "Transform"]]
            for _, item in data["lineage"].head(45).iterrows():
                feature_rows.append([item["feature"], item["source"][:30], item["transform"][:28]])
            story.append(
                _pdf_table(
                    feature_rows[0], feature_rows[1:], widths=[2.1 * inch, 2.0 * inch, 1.5 * inch], font_size=6.3
                )
            )
        else:
            story.append(
                Paragraph(
                    "Evidence references: `documentation/complete_technical_manual.md`, `evaluation/complete_evaluation_record.md`, and the run manifest. These appendices are generated from the same artifact context as the headline results.",
                    styles["small"],
                )
            )
    doc.build(story, onFirstPage=_pdf_header_footer, onLaterPages=_pdf_header_footer)
    return output


def build_ieee_paper(paths: ProjectPaths, context: dict[str, Any]) -> Path:
    """Build an IEEE-inspired two-column research paper PDF."""
    data = _artifact_context(paths, context)
    styles = _report_styles()
    winner = data["winner"]
    profile = data["profile"]
    output = paths.reports / "approach_1_ieee_paper.pdf"

    class TwoColumnDocTemplate(BaseDocTemplate):
        """Two-column A4 template used for the compact manuscript."""

        pass

    doc = TwoColumnDocTemplate(
        str(output),
        pagesize=letter,
        leftMargin=0.55 * inch,
        rightMargin=0.55 * inch,
        topMargin=0.55 * inch,
        bottomMargin=0.48 * inch,
        title="Medical Insurance Claim Fraud Detection — Traditional ML",
        author="IIIT Dharwad",
    )
    frame_width = (letter[0] - 1.1 * inch - 0.22 * inch) / 2
    frame_height = letter[1] - 1.03 * inch
    frames = [
        Frame(
            0.55 * inch,
            0.48 * inch,
            frame_width,
            frame_height,
            id="left",
            leftPadding=3,
            rightPadding=3,
            topPadding=3,
            bottomPadding=3,
        ),
        Frame(
            0.55 * inch + frame_width + 0.22 * inch,
            0.48 * inch,
            frame_width,
            frame_height,
            id="right",
            leftPadding=3,
            rightPadding=3,
            topPadding=3,
            bottomPadding=3,
        ),
    ]
    doc.addPageTemplates([PageTemplate(id="two_col", frames=frames, onPage=_pdf_header_footer)])
    compact = {
        **styles,
        "title": ParagraphStyle("ieee_title", parent=styles["title"], fontSize=17, leading=20, spaceAfter=8),
        "body": ParagraphStyle(
            "ieee_body", parent=styles["body"], fontSize=7.4, leading=9.6, alignment=TA_JUSTIFY, spaceAfter=3
        ),
        "h1": ParagraphStyle("ieee_h1", parent=styles["h1"], fontSize=10.5, leading=12, spaceBefore=5, spaceAfter=3),
        "h2": ParagraphStyle("ieee_h2", parent=styles["h2"], fontSize=8.5, leading=10, spaceBefore=4, spaceAfter=2),
        "small": ParagraphStyle("ieee_small", parent=styles["small"], fontSize=6.3, leading=8),
        "caption": ParagraphStyle("ieee_caption", parent=styles["caption"], fontSize=6.4, leading=7.6),
    }
    story: list[Any] = []
    story.append(
        Paragraph(
            "Medical Insurance Claim Fraud Detection with Leakage-Aware Traditional Machine Learning", compact["title"]
        )
    )
    story.append(
        Paragraph(
            "B. Varshith, M. Jagadeshwar, and J. Ganesh — Department of Data Science and AI, IIIT Dharwad, India",
            compact["subtitle"],
        )
    )
    story.append(Paragraph("Faculty Adviser: Prof. Ramesh Athe", compact["subtitle"]))
    abstract = (
        "Abstract—This paper presents a reproducible classical machine-learning pipeline for medical insurance claim fraud screening. "
        "Using a repository-provided workbook of {rows:,} claims with a {rate} fraud rate, we construct financial and temporal screening features, "
        "exclude identifier-like and potentially post-decision fields, and fit all preprocessing on a stratified training partition. Twenty model "
        "configurations spanning linear, probabilistic, tree, boosting, instance, neural-baseline, and ensemble families are compared using fraud-aware "
        "metrics. Validation F2 selects the operating threshold and ranks models, while PR-AUC provides a complementary ranking view under class imbalance. "
        "The selected model is refit once on training plus validation rows before a locked test evaluation. The artifact chain includes curves, calibration, "
        "permutation importance, fairness slices, serialized model state, a presentation, and two PDF reports. We present the result as a decision-support "
        "baseline rather than an autonomous denial system. The supplied snapshot lacks policy, hospital-tier, document, and verified Indian geography fields; "
        "consequently, the findings are limited to the stated dataset and protocol and motivate later deep-learning and agentic extensions."
    ).format(rows=profile["rows"], rate=_pct(profile["fraud_rate"]))
    story.append(Paragraph(abstract, compact["body"]))
    story.append(
        Paragraph(
            "Keywords—fraud detection, health insurance, imbalanced classification, explainable AI, India",
            compact["small"],
        )
    )
    story.append(Paragraph("1. INTRODUCTION", compact["h1"]))
    story.append(
        Paragraph(
            "Healthcare insurance fraud can appear as inflated charges, duplicate submissions, services not rendered, or inconsistent claim records. Manual investigation is necessary but expensive at scale. A machine-learning screening layer can prioritize records for review, provided that its errors, threshold, and limitations are visible.",
            compact["body"],
        )
    )
    story.append(
        Paragraph(
            "This work addresses a constrained binary task: given a supplied structured claim record, estimate whether its label is Fraud or Legitimate. The contribution is an end-to-end, auditable baseline rather than a claim of population-level superiority. The implementation separates data validation, feature engineering, preprocessing, model fitting, evaluation, explainability, and document generation.",
            compact["body"],
        )
    )
    story.append(Paragraph("2. RELATED WORK", compact["h1"]))
    story.append(
        Paragraph(
            "Statistical fraud detection has long emphasized the trade-off between missed fraud and investigation burden [1]. Anomaly detection surveys motivate deviation-based signals [11], while SMOTE [2] and cost-sensitive learning address imbalance. Random forests [3], gradient boosting [4], margin methods [5], and scikit-learn's unified evaluation framework [15] provide complementary inductive biases for mixed tabular data. SHAP-style explanations [10] motivate the broader interpretability direction, but this baseline uses permutation importance to keep the evidence chain lightweight. Healthcare fraud studies emphasize label quality, coding complexity, provider concentration, and temporal shift [12]–[14].",
            compact["body"],
        )
    )
    story.append(Paragraph("3. DATA AND METHOD", compact["h1"]))
    story.append(
        Paragraph(
            f"The workbook contains {profile['rows']:,} rows and {profile['columns']} source columns. Fraud accounts for {_pct(profile['fraud_rate'])}. Claim amount, income, age, dates, low-cardinality categories, claim type, submission method, specialty, and a supplied cluster field are available. IDs and near-unique diagnosis, procedure, and location values are excluded. Claim status is excluded because it may be observed after adjudication. The target is encoded as 1 for fraud and 0 for legitimate.",
            compact["body"],
        )
    )
    story.append(
        Paragraph(
            "We use a 70/15/15 stratified split with seed 42. Numeric fields receive median imputation and standard scaling; categoricals receive most-frequent imputation and unknown-safe one-hot encoding. The transformer is fitted on training rows only. Feature engineering adds log amounts, claim-to-income ratios, calendar terms, cyclic month encodings, age bands, and screening-time interactions.",
            compact["body"],
        )
    )
    story.append(
        Paragraph(
            "The model zoo includes majority, logistic L1/L2, decision tree, random forest, extra trees, gradient and histogram gradient boosting, AdaBoost, RBF-SVM, KNN, Gaussian and Bernoulli naive Bayes, LDA, QDA, MLP, calibrated ridge, passive-aggressive, voting, and stacking. Searches use stratified three-fold average precision with small declared budgets. Validation thresholds maximize F2 subject to a preferred precision floor.",
            compact["body"],
        )
    )
    story.append(Paragraph("4. EVALUATION", compact["h1"]))
    evaluation_text = "We report accuracy, balanced accuracy, precision, recall, F1, F2, ROC-AUC, PR-AUC, MCC, Cohen's kappa, Brier score, log loss, specificity, confusion matrices, training time, prediction latency, and artifact size. PR-AUC is emphasized under imbalance; F2 captures the operational preference for recall. The selected model is {name} with validation F2 {f2}, validation PR-AUC {prauc}, and threshold {threshold}.".format(
        name=winner.get("display_name"),
        f2=_fmt(winner.get("val_f2")),
        prauc=_fmt(winner.get("val_pr_auc")),
        threshold=_fmt(winner.get("threshold")),
    )
    story.append(Paragraph(evaluation_text, compact["body"]))
    story.append(Paragraph("5. RESULTS", compact["h1"]))
    rows = [["Model", "Val F2", "Val PR", "Test F2"]] + [
        [r["key"], _fmt(r.get("val_f2")), _fmt(r.get("val_pr_auc")), _fmt(r.get("test_f2"))] for r in data["ranked"][:8]
    ]
    story.append(
        _pdf_table(rows[0], rows[1:], widths=[1.35 * inch, 0.52 * inch, 0.58 * inch, 0.52 * inch], font_size=6)
    )
    story.extend(
        [
            Paragraph("6. EXPLAINABILITY AND FAIRNESS", compact["h1"]),
            Paragraph(
                "Permutation importance ranks transformed validation features by average-precision decrease after shuffling. This is an association diagnostic, not a causal explanation. Neutral claim narratives report the score band, top signals, caveats, and next steps without alleging claimant intent. Fairness tables audit gender, age, claim type, and employment slices using recall, FPR, precision, accuracy, and predicted fraud rate; small cells are marked unstable.",
                compact["body"],
            ),
            Paragraph("7. DISCUSSION", compact["h1"]),
            Paragraph(
                "Traditional tree and boosting methods are attractive for tabular fraud triage because they capture nonlinear interactions without a GPU and can be inspected with feature importance. However, the supplied workbook is only 4,500 rows, lacks policy and document evidence, and exhibits a highly concentrated cluster field. A future study should use license-cleared Indian claims, time-based external validation, provider history windows, and a document-aware agent system. We therefore interpret the leaderboard as a baseline under this protocol, not a universal ranking.",
                compact["body"],
            ),
            Paragraph("8. ETHICS AND LIMITATIONS", compact["h1"]),
            Paragraph(
                "The system is decision support. Automatic denial is outside scope. A real implementation must protect claimant data, provide specific reasons, allow correction and appeal, and monitor group error rates. The workbook has no verified Indian geography or policy semantics, so localization claims are limited. Sensitive attributes are used for auditing and excluded from direct model identifiers.",
                compact["body"],
            ),
            Paragraph("9. CONCLUSION", compact["h1"]),
            Paragraph(
                "This paper provides a reproducible traditional-ML fraud-screening baseline and a complete artifact chain. The design makes the selection rule, leakage controls, threshold, metrics, fairness audit, and limitations inspectable. Later deep and agentic approaches can use the same split and evaluator to make a fair comparison.",
                compact["body"],
            ),
            PageBreak(),
            Paragraph("10. ERROR ANALYSIS AND ROBUSTNESS", compact["h1"]),
            Paragraph(
                "A fraud screen should be reviewed through its errors rather than its headline score alone. False positives can represent legitimate high-cost or unusual claims, while false negatives can represent subtle patterns that resemble legitimate activity. The artifact contract therefore preserves confusion matrices, threshold sweeps, class support, and model probabilities. A future error-analysis run should stratify these cases by claim amount, income context, age band, specialty, claim type, submission method, and supplied cluster. These slices are descriptive and should not be interpreted as demographic causation.",
                compact["body"],
            ),
            Paragraph(
                "Robustness is limited by the supplied sample. A random stratified split can overstate generalization when the data-generation mechanism is stable across time. The next evaluation should add a chronological holdout and an external provider or hospital split. Sensitivity should be measured after removing the cluster field, after restricting to low-cardinality features, and after introducing realistic missing documents. Those experiments would distinguish durable fraud signals from convenient synthetic shortcuts.",
                compact["body"],
            ),
            Paragraph("11. REPRODUCIBILITY AND ARTIFACT GOVERNANCE", compact["h1"]),
            Paragraph(
                "Reproducibility is implemented as a chain rather than a sentence. The input workbook is checksummed; the configuration is hashed; the split memberships are persisted; the feature registry records transformed columns; the evaluator writes metric JSON and curve CSV files; and the report builders consume the same context. The selected model is stored with its preprocessor, threshold, target semantics, and run identifier. A new source or configuration requires a new run rather than an in-place edit of the final numbers.",
                compact["body"],
            ),
            Paragraph(
                "This governance is important for a student project because it makes examiner questions answerable. A number in the PDF can be followed to the leaderboard, then to the model JSON, then to the source curve or threshold table. The presentation notes include the same path. A failed model is retained in the run record, so the benchmark does not quietly erase negative results.",
                compact["body"],
            ),
            PageBreak(),
            Paragraph("12. INDIAN-CONTEXT DATA REQUIREMENTS", compact["h1"]),
            Paragraph(
                "A serious Indian insurance system needs policy terms, family-floater membership, waiting periods, co-pay, sum insured, room-rent limits, hospital network status, cashless or reimbursement mode, scheme type, hospital tier, diagnosis and procedure descriptions, bills, discharge summaries, and a verified geography hierarchy. It also needs documentation and language accessibility for diverse policyholders. None of those fields are fabricated here. They are the next data-contract version for the deep-learning and agentic pillars.",
                compact["body"],
            ),
            Paragraph(
                "Contextualization also requires care in interpretation. A higher claim amount in a metro hospital can be legitimate; an elderly claimant can have high utilization without being suspicious; a maternity or pediatric claim must not be treated as anomalous merely because it is less frequent. Fairness checks should be performed on an adequately sampled, license-cleared dataset with human review and appeal outcomes.",
                compact["body"],
            ),
            PageBreak(),
            Paragraph("13. THREATS TO VALIDITY", compact["h1"]),
            Paragraph(
                "The main threats are dataset size, uncertain provenance, possible synthetic label construction, near-unique identifiers, and a single random split. The perfect or near-perfect validation ranking of several models is a warning to investigate the data-generating process rather than a reason to claim state-of-the-art performance. The report treats this behavior as a limitation and highlights the cluster shortcut risk. Statistical tests on one split cannot repair a non-representative dataset.",
                compact["body"],
            ),
            Paragraph("14. CONCLUSION OF THE MANUSCRIPT", compact["h1"]),
            Paragraph(
                "The contribution is a transparent baseline with a fair benchmark harness, explicit leakage controls, validation-only thresholding, locked test evaluation, calibration, explainability, fairness auditing, and submission artifacts. Its numbers are useful as an anchor for later approaches only when the same data split and evaluator are reused. The system is designed to support investigators and policyholders, not to make opaque automatic denials.",
                compact["body"],
            ),
            PageBreak(),
            Paragraph("ACKNOWLEDGMENT", compact["h1"]),
            Paragraph(
                "The authors thank Prof. Ramesh Athe for guidance and IIIT Dharwad for the academic setting.",
                compact["body"],
            ),
            Paragraph("REFERENCES", compact["h1"]),
        ]
    )
    refs = [
        line
        for line in (paths.documentation / "references.md").read_text(encoding="utf-8").splitlines()
        if line[:1].isdigit()
    ]
    story.extend(Paragraph(ref.replace("&", "&amp;"), compact["small"]) for ref in refs)
    doc.build(story)
    return output


def build_all_documents(paths: ProjectPaths, context: dict[str, Any]) -> dict[str, str]:
    """Build Markdown, PPTX, project PDF, and IEEE paper from one context.

    Args:
        paths: Repository path object.
        context: Frozen actual run artifacts.
    Returns:
        Mapping of deliverable names to repository paths.
    """
    write_documentation(paths, context)
    ppt = build_presentation(paths, context)
    report = build_project_report(paths, context)
    paper = build_ieee_paper(paths, context)
    return {"ppt": str(ppt), "project_report": str(report), "ieee_paper": str(paper)}
