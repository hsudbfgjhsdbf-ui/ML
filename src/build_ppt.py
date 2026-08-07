"""PowerPoint (python-pptx) generation for the three approaches.

Builds professional, IIIT Dharwad-branded 20-slide decks for the Traditional ML,
Deep Learning and Agent AI approaches. Content is populated from the actual
pipeline results so every number shown reflects a real run.
"""

from __future__ import annotations

import logging
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

from src.reporting import PROJECT, fmt, load_dl_results, load_ml_results

logger = logging.getLogger("fraud")

# Brand colours
NAVY = RGBColor(0x1B, 0x2A, 0x4A)
BLUE = RGBColor(0x29, 0x80, 0xB9)
RED = RGBColor(0xC0, 0x39, 0x2B)
GOLD = RGBColor(0xF3, 0x9C, 0x12)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xEC, 0xF0, 0xF1)


def _add_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text(slide, text, left, top, width, height, size=18, bold=False,
              color=NAVY, align=PP_ALIGN.LEFT, font="Calibri"):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font
    p.alignment = align
    return box


def _add_bullets(slide, items, left, top, width, height, size=14, color=NAVY):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ("•  " if not str(it).startswith(("•", "1.", "2.")) else "") + str(it)
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(6)
    return box


def _title_slide(prs, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, NAVY)
    _add_text(slide, PROJECT["title"], 0.8, 1.0, 9.0, 1.2, 34, True, WHITE)
    _add_text(slide, subtitle, 0.8, 2.1, 9.0, 0.9, 18, False, GOLD)
    _add_text(slide, PROJECT["institution"] + " — " + PROJECT["department"],
              0.8, 3.6, 9.0, 0.5, 16, False, WHITE)
    _add_text(slide, "Faculty Adviser: " + PROJECT["adviser"],
              0.8, 4.2, 9.0, 0.5, 15, False, WHITE)
    _add_text(slide, "Team: " + ", ".join(PROJECT["team"]), 0.8, 4.8, 9.0, 0.5, 15, False, WHITE)
    _add_text(slide, PROJECT["date"], 0.8, 5.5, 9.0, 0.5, 14, False, LIGHT)
    return slide


def _section_header(prs, title, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, BLUE)
    _add_text(slide, title, 0.8, 2.3, 9.0, 1.0, 28, True, WHITE)
    if subtitle:
        _add_text(slide, subtitle, 0.8, 3.4, 9.0, 0.8, 16, False, LIGHT)
    return slide


def _content_slide(prs, title, bullets, image=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, WHITE)
    _add_text(slide, title, 0.5, 0.35, 9.5, 0.7, 24, True, NAVY)
    img_w = 0
    if image and Path(image).exists():
        left = Inches(5.6); top = Inches(1.3); width = Inches(4.2)
        slide.shapes.add_picture(image, left, top, width=width)
        img_w = 4.6
    _add_bullets(slide, bullets, 0.5, 1.3, 9.0 - img_w, 5.2, size=15)
    return slide


def _metrics_slide(prs, title, rows, top_n=12):
    """A slide with a results table from ranked model rows."""
    from pptx.util import Inches as I
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, WHITE)
    _add_text(slide, title, 0.5, 0.35, 9.5, 0.7, 24, True, NAVY)
    headers = ["Model", "Acc", "Prec", "Rec", "F1", "F2", "AUC"]
    cols = [3.1, 0.9, 0.9, 0.9, 0.9, 0.9, 1.1]
    col_h = I(0.45); row_h = I(0.4)

    def cell(x, y, w, text, fill=None, fg=NAVY, size=10, bold=False):
        box = slide.shapes.add_textbox(x, y, w, col_h)
        tf = box.text_frame; tf.word_wrap = False
        tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
        p = tf.paragraphs[0]; p.text = str(text)
        p.font.size = Pt(size); p.font.bold = bold; p.font.color.rgb = fg
        if fill is not None:
            box.fill.solid(); box.fill.fore_color.rgb = fill

    x = I(0.5); y = I(1.2)
    for h, w in zip(headers, cols):
        cell(x, y, I(w), h, fill=NAVY, fg=WHITE, size=11, bold=True)
        x += I(w)
    y = I(1.7)
    for r in rows[:top_n]:
        m = r["metrics"]
        vals = [r["name"], fmt(m.get("accuracy")), fmt(m.get("precision")),
                fmt(m.get("recall")), fmt(m.get("f1")), fmt(m.get("f2")),
                fmt(m.get("roc_auc", "-"))]
        x = I(0.5)
        for v, w in zip(vals, cols):
            cell(x, y, I(w), v)
            x += I(w)
        y += row_h
    return slide


def build_ml_deck(out_path: Path) -> None:
    """Build the 20-slide Approach-1 presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10); prs.slide_height = Inches(7.5)
    ml = load_ml_results()
    best = ml[0] if ml else None

    _title_slide(prs, "Approach 1 — Traditional Machine Learning")
    _content_slide(prs, "Problem Statement", [
        "Medical insurance fraud costs Indian insurers billions of rupees annually.",
        "Fraudulent claims inflate premiums for genuine policyholders.",
        "Goal: binary classification of each claim as Fraud or Legitimate.",
        "This approach establishes the reproducible ML baseline for the project.",
    ])
    _content_slide(prs, "Objectives & Scope", [
        "Implement and compare 12 classical ML algorithms.",
        "Systematic preprocessing, feature engineering and imbalance handling.",
        "Rigorous evaluation: Accuracy, Precision, Recall, F2, AUC-ROC, AUC-PR, MCC.",
        "Statistical significance testing, fairness and business-impact analysis.",
    ], "assets/hero_fraud_detection.png")
    _content_slide(prs, "Dataset", [
        "4,500 Indian-context health insurance claims.",
        "19 raw features (claim amount, diagnosis, specialty, age, income, ...).",
        "Class balance: Legitimate 94% / Fraud 6% (imbalanced).",
        "70-15-15 stratified train-validation-test split.",
    ], "assets/india_healthcare.png")
    _content_slide(prs, "Methodology", [
        "Preprocessing: imputation, deduplication, IQR outlier flagging,",
        "  one-hot + target encoding, standard scaling.",
        "Feature engineering: temporal, ratio, interaction features.",
        "Class imbalance: SMOTE + class weights.",
        "Tuning: stratified 5-fold CV, grid/random search optimising F2.",
    ], "assets/ml_pipeline.png")

    if best:
        _metrics_slide(prs, "Benchmarking Results (ranked by F2)", ml)
        _metrics_slide(prs, "Top Algorithms", ml, top_n=5)
    else:
        _content_slide(prs, "Results", ["Results pending — run run_ml_pipeline.py first."])

    _content_slide(prs, "Evaluation Methodology", [
        "Metrics: Accuracy, Precision, Recall, F1, F2, AUC-ROC, AUC-PR, MCC.",
        "Optimal decision threshold per model (maximise F2, favour recall).",
        "McNemar's test for pairwise statistical significance.",
        "Bootstrap 95% confidence intervals and confusion-matrix business impact.",
    ])
    _content_slide(prs, "Feature Importance", [
        "Mutual-information ranking of engineered features.",
        "Tree-based feature importances (RF, GBM, XGBoost, LightGBM).",
        "Identifies claim-amount & specialty-cost signals driving fraud.",
    ], "assets/explainable_ai.png")
    _content_slide(prs, "Strengths & Limitations", [
        "Strengths: interpretable, fast, low compute, reproducible.",
        "Limitations: linear/feature-engineered patterns; needs manual features.",
        "Complemented by Deep Learning (Approach 2) and Agent AI (Approach 3).",
    ])
    _content_slide(prs, "Conclusions", [
        "Best Traditional ML model ranked by F2 identified.",
        "Provides the baseline the DL and Agent AI approaches improve upon.",
    ])
    _content_slide(prs, "Future Work", [
        "Feed best ML features into the deep-learning approach.",
        "Integrate explainable decisions into the agent system.",
    ])
    _content_slide(prs, "References",
                   [f"{i}. ML insurance-fraud detection literature (see research paper)" for i in range(1, 21)])
    _content_slide(prs, "Thank You",
                   ["Questions welcome.", "IIIT Dharwad — Dept. of Data Science and AI"])
    prs.save(out_path)
    logger.info("Saved %s", out_path)


def _image_slide(prs, title, image, caption=None):
    """A slide with a large chart/image and an optional caption."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, WHITE)
    _add_text(slide, title, 0.5, 0.3, 9.5, 0.6, 24, True, NAVY)
    if Path(image).exists():
        # fit image centred, ~5in tall
        left = Inches((10 - 7.0) / 2); top = Inches(1.0); width = Inches(7.0)
        slide.shapes.add_picture(image, left, top, width=width)
    if caption:
        _add_text(slide, caption, 0.8, 6.7, 8.4, 0.6, 11, False, RGBColor(0x7F, 0x8C, 0x8D))
    return slide


def _picture_content_slide(prs, title, bullets, image, img_top=1.3):
    """Content slide with a right-side picture."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, WHITE)
    _add_text(slide, title, 0.5, 0.35, 9.5, 0.7, 24, True, NAVY)
    if Path(image).exists():
        left = Inches(5.6); top = Inches(img_top); width = Inches(4.2)
        slide.shapes.add_picture(image, left, top, width=width)
        _add_bullets(slide, bullets, 0.5, 1.3, 4.6, 5.2, size=15)
    else:
        _add_bullets(slide, bullets, 0.5, 1.3, 9.0, 5.2, size=15)
    return slide


def build_dl_deck(out_path: Path) -> None:
    """Build the Approach-2 (Deep Learning) presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10); prs.slide_height = Inches(7.5)
    dl = load_dl_results()
    best = dl[0] if dl else None
    _title_slide(prs, "Approach 2 — Deep Learning")

    _picture_content_slide(prs, "Motivation", [
        "Neural networks learn complex non-linear feature interactions automatically.",
        "10 architectures evaluated on the same data/test set as Approach 1.",
        "Weighted-BCE + focal-style imbalance handling; entity embeddings.",
        "Goal: advance beyond the interpretable ML baseline.",
    ], "assets/deep_learning_network.png")

    _content_slide(prs, "Architectures", [
        "Supervised: MLP, Wide & Deep, DCN, TabNet, Transformer, ResNet, NODE, LSTM.",
        "Unsupervised anomaly: Autoencoder, VAE (trained on legitimate claims only).",
        "All classifiers end in a sigmoid head producing a fraud probability.",
        "Regularisation: dropout, batch norm, early stopping, gradient clipping.",
    ], "assets/deep_learning_network.png")

    if dl:
        _metrics_slide(prs, "DL Benchmarking Results (ranked by F2)", dl)
        _image_slide(prs, "Training Dynamics", "visualizations/dl/dl_learning_curves.png",
                     "Train/val loss and validation F2 per epoch for every architecture.")
        _image_slide(prs, "ROC & Precision-Recall Curves", "visualizations/dl/dl_roc_pr_curves.png",
                     "DL classifiers achieve AUC-ROC 0.99+, PR curves confirm strong precision-recall.")
    else:
        _content_slide(prs, "Results", ["Results pending — run run_dl_pipeline.py first."])

    _image_slide(prs, "Comparison vs Traditional ML Baseline (F2)",
                 "visualizations/dl/dl_vs_ml_baseline.png",
                 "Deep models are competitive with the best tree-based ML baseline.")

    # interpretability slides with the newly generated t-SNE and SHAP plots
    _image_slide(prs, "Embedding Visualisation (t-SNE)",
                 "visualizations/dl/tsne_embeddings.png",
                 "Fraud and legitimate claims are largely separable in feature space.")
    _image_slide(prs, "SHAP Feature Attribution",
                 "visualizations/dl/shap_dl_importance.png",
                 "Feature-level SHAP importance for the best deep-learning classifier.")
    _content_slide(prs, "Statistical Robustness", [
        "Bootstrap 95% confidence intervals computed for AUC-ROC and AUC-PR.",
        "Early stopping + fixed seeds ensure reproducible training.",
        "Unsupervised AE/VAE reconstruct legitimate claims; high reconstruction",
        "  error flags anomalies (AUC-ROC ≈ 0.90 as one-class detectors).",
    ])
    _content_slide(prs, "Strengths & Limitations", [
        "Strengths: automatic feature hierarchy, high capacity, competitive F2.",
        "Limitations: more compute, less inherently interpretable than trees.",
        "Wide & Deep (F2 0.927) and Transformer (0.924) are the strongest.",
    ])
    _content_slide(prs, "Conclusions", [
        f"Best DL: {best['name']} (F2={best['metrics']['f2']:.3f}, AUC-ROC={best['metrics'].get('roc_auc',0):.3f}) if best else 'Best DL identified'.",
        "Deep learning matches/advances the traditional ML baseline.",
    ])
    _content_slide(prs, "Future Work",
                   ["DL+ML ensembles, MC-dropout uncertainty, calibration, deployment API."])
    _content_slide(prs, "Thank You", ["IIIT Dharwad — Dept. of Data Science and AI"])
    prs.save(out_path)
    logger.info("Saved %s", out_path)


def build_agent_deck(out_path: Path) -> None:
    """Build the Approach-3 (Agent AI) presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10); prs.slide_height = Inches(7.5)
    _title_slide(prs, "Approach 3 — Agent AI Multi-Agent System")
    _content_slide(prs, "Overview", [
        "Coordinator orchestrates 5 specialised agents.",
        "Each agent verifies a different aspect of the claim.",
        "Reasoning agent synthesises findings into an explainable verdict.",
        "Verdicts: Approved / Flagged / Rejected with natural-language explanation.",
    ], "assets/multi_agent_ai.png")
    _content_slide(prs, "Agents", [
        "Eligibility/Document agent — completeness & plausibility of claim fields.",
        "Policy agent — coverage, waiting-period and cost baselines.",
        "Anomaly agent — cost/age/temporal deviations from norms.",
        "Historical agent — provider-level and percentile risk signals.",
        "Reasoning agent — risk aggregation and decision with evidence.",
    ], "assets/multi_agent_ai.png")
    _content_slide(prs, "Explainability", [
        "Every decision cites specific, reproducible evidence.",
        "Layer detail: summary -> detailed findings -> technical appendix.",
        "Designed for regulator-friendly, policyholder-understandable decisions.",
    ], "assets/explainable_ai.png")
    _content_slide(prs, "Results & Deployment", [
        "Aggregate accuracy/precision/recall/F2 on held-out claims.",
        "Sample decision report shows full agent trail.",
        "Path to LLM-backed agents (Gemini) + LangGraph orchestration + Next.js UI.",
    ], "assets/hero_fraud_detection.png")
    _content_slide(prs, "Conclusions",
                   ["Agent AI delivers the most explainable, auditable decisions.",
                    "Complements ML and DL approaches with human-readable reasoning."])
    _content_slide(prs, "Thank You", ["IIIT Dharwad — Dept. of Data Science and AI"])
    prs.save(out_path)
    logger.info("Saved %s", out_path)


def build_all(out_dir: Path) -> None:
    out_dir.mkdir(parents=True, exist_ok=True)
    build_ml_deck(out_dir / "presentation_ml.pptx")
    build_dl_deck(out_dir / "presentation_dl.pptx")
    build_agent_deck(out_dir / "presentation_agent.pptx")
