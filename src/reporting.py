"""Generate traceable Approach-1 Markdown, IEEE-style PDF and 20-slide presentation.

Narrative files are rendered from the current evaluated benchmark rather than hand-filled
numbers. This makes the result claims, charts, deck and report internally consistent.
"""
from __future__ import annotations

from datetime import datetime, timezone
from pathlib import Path
from typing import Any

import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import cm
from reportlab.platypus import BaseDocTemplate, Frame, Image, KeepTogether, PageBreak, Paragraph, Spacer, Table, TableStyle

from src.feature_engineering import feature_dictionary

REFERENCES = [
    "[1] C. Phua, V. Lee, K. Smith and R. Gayler, “A comprehensive survey of data mining-based fraud detection research,” Artificial Intelligence Review, vol. 34, pp. 1–14, 2010.",
    "[2] E. W. T. Ngai, Y. Hu, Y. H. Wong, Y. Chen and X. Sun, “The application of data mining techniques in financial fraud detection,” Decision Support Systems, vol. 50, no. 3, pp. 559–569, 2011.",
    "[3] R. J. Bolton and D. J. Hand, “Statistical fraud detection: A review,” Statistical Science, vol. 17, no. 3, pp. 235–255, 2002.",
    "[4] A. Abdallah, M. A. Maarof and A. Zainal, “Fraud detection system: A survey,” Journal of Network and Computer Applications, vol. 68, pp. 90–113, 2016.",
    "[5] R. A. Bauder and T. M. Khoshgoftaar, “The detection of Medicare fraud using machine learning methods with excluded provider labels,” Health Care Management Science, vol. 20, pp. 1–15, 2017.",
    "[6] H. Joudaki et al., “Using data mining to detect health care fraud and abuse: A review of literature,” Global Journal of Health Science, vol. 7, no. 1, pp. 194–202, 2015.",
    "[7] N. Japkowicz and S. Stephen, “The class imbalance problem: A systematic study,” Intelligent Data Analysis, vol. 6, no. 5, pp. 429–449, 2002.",
    "[8] N. V. Chawla, K. W. Bowyer, L. O. Hall and W. P. Kegelmeyer, “SMOTE: Synthetic minority over-sampling technique,” Journal of Artificial Intelligence Research, vol. 16, pp. 321–357, 2002.",
    "[9] H. He and E. A. Garcia, “Learning from imbalanced data,” IEEE Transactions on Knowledge and Data Engineering, vol. 21, no. 9, pp. 1263–1284, 2009.",
    "[10] L. Breiman, “Random forests,” Machine Learning, vol. 45, no. 1, pp. 5–32, 2001.",
    "[11] J. H. Friedman, “Greedy function approximation: A gradient boosting machine,” Annals of Statistics, vol. 29, no. 5, pp. 1189–1232, 2001.",
    "[12] T. Chen and C. Guestrin, “XGBoost: A scalable tree boosting system,” in Proc. KDD, 2016, pp. 785–794.",
    "[13] G. Ke et al., “LightGBM: A highly efficient gradient boosting decision tree,” in Proc. NeurIPS, 2017, pp. 3146–3154.",
    "[14] C. Cortes and V. Vapnik, “Support-vector networks,” Machine Learning, vol. 20, pp. 273–297, 1995.",
    "[15] Y. Freund and R. E. Schapire, “A decision-theoretic generalization of on-line learning and an application to boosting,” Journal of Computer and System Sciences, vol. 55, no. 1, pp. 119–139, 1997.",
    "[16] L. V. Utkin, “A method for processing imprecise expert judgments about parameters of probability distributions,” European Journal of Operational Research, vol. 158, no. 3, pp. 657–674, 2004.",
    "[17] S. M. Lundberg and S.-I. Lee, “A unified approach to interpreting model predictions,” in Proc. NeurIPS, 2017, pp. 4765–4774.",
    "[18] J. Davis and M. Goadrich, “The relationship between Precision-Recall and ROC curves,” in Proc. ICML, 2006, pp. 233–240.",
    "[19] M. P. Wand and M. C. Jones, Kernel Smoothing. London, U.K.: Chapman & Hall, 1995.",
    "[20] M. Hardt, E. Price and N. Srebro, “Equality of opportunity in supervised learning,” in Proc. NeurIPS, 2016, pp. 3315–3323.",
    "[21] F. Kamiran and T. Calders, “Data preprocessing techniques for classification without discrimination,” Knowledge and Information Systems, vol. 33, pp. 1–33, 2012.",
    "[22] Q. McNemar, “Note on the sampling error of the difference between correlated proportions,” Psychometrika, vol. 12, pp. 153–157, 1947.",
    "[23] Insurance Regulatory and Development Authority of India, “Protection of Policyholders’ Interests Regulations,” IRDAI, New Delhi, India, 2024.",
]


def _fmt(value: Any, digits: int = 3) -> str:
    """Render metric values safely for Markdown, slides and tables."""
    if isinstance(value, (float, np.floating)):
        return "—" if not np.isfinite(value) else f"{value:.{digits}f}"
    return str(value)


def _md_table(frame: pd.DataFrame, columns: list[str], rename: dict[str, str] | None = None, max_rows: int | None = None) -> list[str]:
    """Convert selected DataFrame columns into a readable Markdown table."""
    rename = rename or {}
    subset = frame.loc[:, columns].copy()
    if max_rows:
        subset = subset.head(max_rows)
    headers = [rename.get(c, c) for c in columns]
    rows = ["| " + " | ".join(headers) + " |", "|" + "|".join(["---"]*len(headers)) + "|"]
    for _, record in subset.iterrows():
        rows.append("| " + " | ".join(_fmt(record[c]).replace("|", "/") for c in columns) + " |")
    return rows


def _asset(root: Path, category: str, filename: str) -> str:
    """Return repository-relative image path for reliable Markdown embedding."""
    path = root / category / filename
    return str(path.relative_to(root.parents[0])) if path.exists() else str(path)


def generate_markdown(config: dict[str, Any], source_audit: dict[str, Any], quality: dict[str, Any], benchmark: pd.DataFrame, fairness: pd.DataFrame, imbalance: pd.DataFrame, significance: pd.DataFrame) -> tuple[Path, Path]:
    """Write complete Approach-1 documentation and evaluation reports in single writes.

    Args:
        config: Loaded configuration.
        source_audit: Original workbook adequacy audit.
        quality: Cleaning/split/feature-selection evidence.
        benchmark: Full benchmark table with private prediction columns.
        fairness: Best-model group audit.
        imbalance: Validation sampling-strategy comparison.
        significance: Pairwise statistical tests.

    Returns:
        Paths to the documentation and evaluation Markdown files.
    """
    docs = Path(config["paths"]["documentation_dir"]); evaluation = Path(config["paths"]["evaluation_dir"])
    docs.mkdir(parents=True, exist_ok=True); evaluation.mkdir(parents=True, exist_ok=True)
    public = benchmark.drop(columns=[c for c in benchmark.columns if c.startswith("_")], errors="ignore").sort_values(["f2","auc_roc"], ascending=False).reset_index(drop=True)
    best = public.iloc[0]
    visual_root = Path(config["paths"]["visualization_dir"])
    now = datetime.now(timezone.utc).strftime("%Y-%m-%d %H:%M UTC")
    lines: list[str] = [
        "# Approach 1 — Traditional Machine Learning for Indian Medical Insurance Claim Fraud Screening",
        "", "**Institution:** IIIT Dharwad, Department of Data Science and AI  ", "**Faculty Adviser:** Ramesh Athe  ",
        "**Team:** B Varshith (23BDS011), M Jagadeshwar (23BDS033), J Ganesh (23BDS024)  ",
        f"**Generated from the verified run:** {now}  ", "**Data status:** Transparent synthetic Indian-context educational data; not an insurer dataset and not a real-claim decision system.", "",
        "> **Safety statement.** The output is a fraud-screening recommendation for human investigation. A high score must never autonomously deny an insurance claim. Coverage, medical necessity, claimant rights and a human reviewer remain decisive.", "",
        "## Abstract", "",
        f"This study implements a reproducible traditional-machine-learning baseline for medical insurance claim fraud screening in an Indian context. The supplied workbook was retained unchanged but did not meet the study adequacy criteria: it contains {source_audit['rows']:,} records, generic rather than Indian locations, and no required policy/waiting-period/history features. We therefore generated a documented synthetic population, preserved realistic Indian product categories and INR price variation, and used stratified 70/15/15 partitions. A train-only pipeline applies duplicate auditing, median/mode imputation, robust scaling, one-hot encoding, smoothed target encoding for high-cardinality fields, limited domain interactions and mutual-information feature selection. The benchmark evaluates {len(public)} classifiers using F2 as the selection metric because missed fraud is costly. The current held-out leader is **{best['algorithm']}** (F2={best['f2']:.3f}, recall={best['recall']:.3f}, PR-AUC={best['auc_pr']:.3f}). Results are pedagogical and describe synthetic behavior only; they do not establish expected performance at an Indian insurer. The report includes INR-oriented cost assumptions, threshold selection, statistical comparison, explanation assets and group-level fairness audit.", "",
        "**Keywords—** health-insurance fraud, India, imbalanced classification, explainable machine learning, fairness audit, F2 score, synthetic data.", "",
        "## 1. Problem, Scope and Research Questions", "",
        "Fraudulent bills, duplicated claims, inflated procedures and policy-timing manipulation can impose losses on insurers and ultimately raise premiums for genuine policyholders. India adds material context: family-floater covers, employer-group policies, Ayushman Bharat and ECHS schemes, allopathic and Ayurvedic treatments, cashless and reimbursement workflows, GST components and strong tier/region cost variation. A ₹2 lakh procedure at a metro corporate hospital cannot be evaluated with the same baseline as a small-town nursing-home episode.", "",
        "This approach answers four bounded questions: (1) which classical model has the strongest recall-prioritised performance on the defined synthetic test set; (2) what threshold best implements the chosen F2 objective; (3) which observed synthetic signals contribute to model screening; and (4) whether error rates differ materially across audited groups. The study does not infer actual national fraud prevalence, clinical validity or insurer-specific policy eligibility.", "",
        "### Objectives", "",
        "1. Build an inspectable end-to-end binary fraud-screening pipeline.", "2. Preserve and audit the bundled workbook while transparently using a documented fallback population when it is inadequate.", "3. Evaluate multiple model families under the same split and F2-led protocol.", "4. Report threshold, probability quality, INR assumptions, computational cost, explanation and fairness evidence.", "5. Produce reusable assets for the presentation and IEEE-style report.", "",
        "## 2. Data Provenance, Adequacy and Ethics", "",
        "### 2.1 Bundled-workbook audit", "",
        f"The supplied `Health Insurance Fraud Claims.xlsx` has SHA-256 `{source_audit['sha256']}` and {source_audit['rows']:,} data rows. It includes a binary `ClaimLegitimacy` field, but it falls below the 10,000-record criterion and lacks Indian policy/product fields, temporal policy controls and reliable claim-history features. The source was copied unchanged to `data/raw/`; it was not relabelled, currency-converted or represented as an Indian insurer extract.", "",
        "### 2.2 Synthetic fallback design", "",
        "The generated population contains claimants across Indian states/cities, INR amounts with a log-normal/right-skewed tail, insurer/policy types, government schemes, hospital tiers, coverage and waiting-period fields, allopathic/Ayurvedic treatment types, historical features and a noisy fraud-generating mechanism. It contains deliberately duplicated records solely to test duplicate removal. Fraud probability is probabilistic and has unobserved noise, so no feature is a deterministic label key. The complete generator and fixed seed are in `src/data_loading.py` and `configs/traditional_ml.yaml`.", "",
        "### 2.3 Privacy and intended use", "",
        "All values are synthetic. No Aadhaar number, PAN, real hospital, real patient, medical image or API secret is included. Demographic variables are retained to audit harms, rather than to define an automatic adverse action. Any real deployment would require lawful data governance, insurer-specific validation, privacy impact assessment, security review, clinical/policy expert review, calibration monitoring and meaningful human appeal.", "",
        "## 3. Literature Review and Identified Gap", "",
        "The following work motivates the model families and safeguards. It is not treated as evidence that a synthetic benchmark transfers to live claims.", "",
        "| Ref. | Contribution | Dataset/method | Finding relevant here | Limitation addressed in this project |", "|---|---|---|---|---|",
        "| [1] | Broad fraud-mining survey | Multiple fraud domains | Feature engineering and imbalanced evaluation matter | Does not specialise to Indian medical policies |",
        "| [2] | Financial-fraud data-mining review | Classification and clustering | Supervised learning is valuable when labels exist | Domain context and fairness are often underdeveloped |",
        "| [3] | Statistical fraud review | Behavioural anomaly methods | Fraud evolves and distributions matter | This baseline is explicitly bounded to a fixed synthetic period |",
        "| [4] | Fraud-system survey | Rule/ML approaches | Hybrid evidence is useful | Motivates later evidence-grounded Agent-AI track |",
        "| [5] | Medicare fraud detection | Provider-label learning | Provider patterns can be informative | Current-provider proxies are disclosed synthetic variables |",
        "| [6] | Health-care fraud review | Health-data mining | Interpretability and validation are essential | Adds threshold, cost and fairness audit |",
        "| [7] | Class imbalance study | Imbalanced supervised learning | Accuracy can mislead | F2/PR-AUC/MCC are primary reports |",
        "| [8] | SMOTE | Minority synthesis | Sampling can improve minority learning | Sampling comparison is train-only and documented |",
        "| [9] | Imbalanced-learning review | Survey | Cost-sensitive choices require context | INR assumptions are explicit and illustrative |",
        "| [10] | Random forests | Bagged trees | Nonlinear tabular patterns and importances | Importances are not causal explanations |",
        "| [11] | Gradient boosting | Functional gradient descent | Boosting is strong on structured data | Tuned under common split/metric |",
        "| [12] | XGBoost | Regularised tree boosting | Efficient high-performance trees | Optional dependency is recorded, not assumed |",
        "| [13] | LightGBM | Leaf-wise boosting | Speed/accuracy trade-off | Latency/model size are benchmarked |",
        "| [14] | Support-vector networks | Margin classifiers | Kernels model nonlinear boundaries | SVM training cap is disclosed for laptop reproducibility |",
        "| [15] | AdaBoost | Adaptive ensembles | Weak learners can combine effectively | Performance verified on held-out data |",
        "| [17] | SHAP | Additive attributions | Local/global model explanation framework | Attribution is explained as association, not evidence of fraud |",
        "| [18] | PR vs ROC | Evaluation theory | PR better reveals positive-class performance | Both PR-AUC and ROC-AUC are retained |",
        "| [20] | Equality of opportunity | Fair classification | FPR/FNR differences matter | Group outcomes are explicitly audited |",
        "| [21] | Discrimination-aware preprocessing | Fair ML | Mitigation has trade-offs | Guardrails and mitigation route are documented |",
        "| [22] | McNemar test | Paired classifiers | Same-test-set predictions should be compared paired | Exact test is reported descriptively |",
        "",
        "**Gap.** Existing general fraud studies do not by themselves yield a transparent, reproducible, Indian-context claim-screening baseline with policy timing, regional costs, sampling comparison, threshold selection, INR-risk assumptions and demographic error audit. This educational synthetic implementation provides that engineering baseline; it does not fill the need for a real governed claims dataset.", "",
        "## 4. Methodology", "",
        f"![Traditional ML pipeline]({_asset(visual_root, 'technical', 'traditional_ml_pipeline_diagram.png')})", "",
        "### 4.1 Reproducible split and leakage controls", "",
        f"Exact duplicate claims were removed before the stratified split. The resulting partitions are train={quality['stratified_split']['train_rows']:,}, validation={quality['stratified_split']['validation_rows']:,}, test={quality['stratified_split']['test_rows']:,}; fraud rates are {quality['stratified_split']['train_fraud_rate']:.2%}, {quality['stratified_split']['validation_fraud_rate']:.2%} and {quality['stratified_split']['test_fraud_rate']:.2%}. The preprocessor, target encoder, scaler and mutual-information selector are fitted only on training rows. Validation chooses threshold; test is used only for final model comparison.", "",
        "### 4.2 Cleaning and outlier policy", "",
        f"Missingness was measured before imputation. Exact duplicates removed: {quality['exact_duplicates_removed']}. Numeric values use median imputation; low-cardinality categoricals use the training mode; target encoding maps unknown high-cardinality categories to the training global rate. Z-score and IQR candidates are reported but preserved when plausible because a very large claim can be precisely the signal under investigation. Clearly impossible values would be rejected by schema validation; none are silently corrected.", "",
        "### 4.3 Encoding and scaling", "",
        f"The fitted input has {quality['feature_selection']['features_before_selection']} transformed features: numeric fields use robust scaling, low-cardinality nominal fields use one-hot encoding and configured/high-cardinality fields use smoothed training-only target encoding. Robust scaling reduces leverage from INR tails for logistic, SVM, KNN, QDA and neural-network bridge models. Multinomial Naive Bayes additionally receives a fitted train-only min–max transformation because its likelihood requires nonnegative values.", "",
        "### 4.4 Domain features", "",
        "The feature engineering code creates only claim-time functions: claim-to-premium ratio, regional treatment-cost deviation, days relative to waiting-period end, claim-frequency intensity, amount per hospital day, current-to-history ratio, provider-distance interaction, age×amount, policy utilisation and two limited degree-two transforms. The regional baseline normalises tier/city cost rather than treating metro pricing as inherently suspicious.", "",
        "| Engineered feature | Definition | Why it may be useful | Caution |", "|---|---|---|---|",
    ]
    for item in feature_dictionary():
        lines.append(f"| `{item['feature']}` | {item['description']} | {item['fraud_relevance']} | Association is not proof; a reviewer must verify evidence. |")
    lines += ["", "### 4.5 Feature selection", "", f"Correlation inspection identifies redundant numerical relationships. Mutual-information filtering selected {quality['feature_selection']['features_after_selection']} training-fitted columns, listed in `data/processed/data_quality_report.json`. RFE/LASSO/tree importance remain complementary analytical techniques; this shared filter controls dimensionality fairly across the complete benchmark rather than optimizing separately on the test set.", "", "### 4.6 Imbalance strategies", "", "Five strategies were compared using a train-only logistic reference and validation-only threshold selection: class weighting, random undersampling, Tomek links, SMOTE and SMOTEENN. This does not declare that a single sampler is globally optimal; it documents the recall/precision trade-off prior to the all-model class-weight baseline.", ""]
    lines += _md_table(imbalance, ["strategy","training_rows_after_sampling","threshold","validation_f2","validation_recall","validation_precision"], {"training_rows_after_sampling":"train rows","validation_f2":"val F2","validation_recall":"val recall","validation_precision":"val precision"})
    lines += ["", "### 4.7 Algorithms and tuning protocol", "", "All models use five-fold stratified cross-validation on the training partition and select hyperparameters by mean F2. Small grids are exhaustive; broader tree/neural spaces use bounded random search. The selected threshold is not fixed at 0.50: validation probabilities are evaluated from 0.05 to 0.95 in 0.005 steps and the highest F2 is retained. F2 weights recall four times as strongly as precision in its denominator, reflecting the missed-fraud screening priority while retaining precision, cost and FPR reporting.", ""]
    for _, row in public.iterrows():
        lines += [f"#### {row['algorithm']}", f"Search space: `{row['best_hyperparameters']}` selected from `{row.get('search_space', 'see saved model metadata')}`. Best five-fold CV F2={row['cv_f2_mean']:.3f}±{row['cv_f2_std']:.3f}; validation threshold={row['threshold']:.3f}; tuning time={row['training_time_seconds']:.2f}s. The model is serialized with metadata in `models/`.", ""]
    lines += ["## 5. Experimental Results", "", "### 5.1 Held-out benchmark", "", "The following values are on the untouched test partition, at the model-specific threshold chosen from validation F2. The primary rank is F2 and the secondary rank is ROC-AUC. Accuracy is intentionally not the ranking metric because the synthetic class is imbalanced.", ""]
    lines += _md_table(public, ["algorithm","accuracy","precision","recall","f1","f2","auc_roc","auc_pr","mcc","threshold","training_time_seconds","prediction_time_per_sample_ms","model_size_kb","tuned_hyperparameters"], {"training_time_seconds":"train s","prediction_time_per_sample_ms":"latency ms/sample","tuned_hyperparameters":"# tuned"})
    lines += ["", f"**Best screened model:** `{best['algorithm']}`. Its test F2={best['f2']:.3f}, recall={best['recall']:.3f}, precision={best['precision']:.3f}, PR-AUC={best['auc_pr']:.3f}, ROC-AUC={best['auc_roc']:.3f} and MCC={best['mcc']:.3f}. The threshold ({best['threshold']:.3f}) was selected without inspecting test labels.", "", f"![Grouped metrics]({_asset(visual_root, 'model_comparison', 'grouped_metric_comparison.png')})", "", f"![ROC curves]({_asset(visual_root, 'model_comparison', 'roc_curves_all_models.png')})", "", f"![Precision recall curves]({_asset(visual_root, 'model_comparison', 'precision_recall_curves_all_models.png')})", "", f"![Efficiency]({_asset(visual_root, 'model_comparison', 'training_time_vs_accuracy.png')})", "", "### 5.2 Cost-sensitive confusion analysis", "", "A false negative means a synthetic fraudulent claim was recommended for approval; its illustrative loss is that claim’s INR amount. A false positive is a legitimate claim routed to review, charged at the configured ₹3,500 review/friction proxy. These are scenario assumptions, not actual insurer costs, recoveries or regulatory exposure.", ""]
    cost_rows = public[["algorithm","false_negative","false_positive","cost_matrix_inr"]].copy()
    for _, row in cost_rows.iterrows():
        cost = row["cost_matrix_inr"] if isinstance(row["cost_matrix_inr"], dict) else {}
        lines.append(f"- **{row['algorithm']}:** FN={row['false_negative']}, FP={row['false_positive']}, illustrative total cost=₹{float(cost.get('total_illustrative_cost_inr', float('nan'))):,.2f}.")
    lines += ["", f"![Confusion matrices]({_asset(visual_root, 'model_comparison', 'confusion_matrices_all_models.png')})", "", "### 5.3 Paired significance evidence", "", "McNemar’s exact test compares paired correctness on the same test records; Wilcoxon compares five cross-validation F2 fold values. p<0.05 is a descriptive threshold only: synthetic records are generated rather than independent insurer observations, and multiple comparisons increase false-positive risk.", ""]
    if not significance.empty:
        lines += _md_table(significance, list(significance.columns))
    else:
        lines.append("Only one benchmark row was available; pairwise tests were not applicable.")
    lines += ["", "## 6. Explainability and Error Analysis", "", f"Tree split gain/importances, random-forest importances and logistic coefficient magnitudes are rendered in `visualizations/interpretability/`. These methods identify predictive association in the synthetic study, not causal evidence that a person or provider committed fraud. For a claimant-facing explanation, the system should name verifiable evidence (e.g., policy waiting-period status, itemised bill mismatch) and offer a human review path rather than expose opaque raw scores.", "", "The synthetic error analysis should examine false negatives with very high regional treatment costs, claims near waiting-period end and incomplete documents; these may represent legitimate exceptional care. False positives may also concentrate in expensive metro/corporate treatment episodes. The correct operational response is manual evidence review, not automatic refusal.", "", "## 7. Fairness and Ethical Analysis", "", "Group metrics are calculated for the F2-leading model across gender, age bracket, state, income bracket and treatment type. A difference above the configured 0.10 guardrail in FPR/FNR requires investigation, sample-size checks and a mitigation comparison (for example, reweighting or removal of a proxy feature). Small groups are explicitly marked and should not support strong conclusions.", ""]
    fair_summary = fairness.groupby("dimension").agg(n_groups=("group","count"), max_fpr_gap=("fpr",lambda s: float(s.max()-s.min())), max_fnr_gap=("fnr",lambda s: float(s.max()-s.min())), accuracy_min=("accuracy","min"), accuracy_max=("accuracy","max")).reset_index()
    lines += _md_table(fair_summary, list(fair_summary.columns))
    lines += ["", f"![Fairness FNR audit]({_asset(visual_root, 'fairness', 'fairness_fnr_groups.png')})", "", f"![Calibration]({_asset(visual_root, 'fairness', 'calibration_reliability_diagram.png')})", "", "The audit supports monitoring rather than demographic rationing. Gender, disability, age and geography must not be used as shortcut reasons to reject a claim. In a production setting, protected attributes should be access-controlled and separated from scoring unless legally/ethically justified for a carefully governed fairness intervention.", "", "## 8. Code Walkthrough and Reproduction", "", "| File | Responsibility | Key outputs / checks |", "|---|---|---|", "| `src/data_loading.py` | audits the raw workbook and generates the labelled fallback population | checksum, adequacy report, synthetic data metadata |", "| `src/preprocessing.py` | duplicate/outlier audit, train-only encoding/scaling/selection and splits | serialized preprocessor/selector, processed partitions |", "| `src/feature_engineering.py` | deterministic domain/temporal/interaction features | documented claim-time variables |", "| `src/models.py` | classifier registry, 5-fold F2 tuning and validation threshold selection | saved models and hyperparameters |", "| `src/evaluate.py` | held-out metrics, INR proxy costs, fairness, calibration, significance | CSV/JSON evaluation evidence |", "| `src/visualize.py` | EDA, comparison, explanation and audit plots | reusable PNG assets |", "| `src/reporting.py` | derives Markdown, deck and report from live outputs | internally consistent academic artefacts |", "| `src/train.py` | one-command orchestrator | end-to-end verification manifest |", "", "Run from a fresh virtual environment: `python -m venv .venv && .venv/bin/pip install -r requirements.txt && .venv/bin/python -m src.train --regenerate-data`. The fixed seed, package versions, configurations and source checksum are written alongside the artifacts. A clean rerun overwrites generated output with the new verified run; never mix files across configurations.", "", "## 9. Limitations and Future Work", "", "1. Labels, price distributions and fraud mechanisms are synthetic, so metric values are not estimates of operational accuracy or financial benefit.\n2. Random claim-level splitting can retain related policyholders/providers across partitions; future real-data work should add temporal and entity-disjoint validation.\n3. Target encoding can encode historical correlation, not causation; it must be monitored for provider/demographic proxy harm.\n4. Synthetic historical/provider rates are simplified and must not be conflated with verified investigations.\n5. Classifier explanations do not authenticate bills or establish legal fraud; Approach 3 is designed to add document evidence, RAG citations and human checkpoints.\n6. Deployment needs calibration drift monitoring, independent audit, policy-specific coverage rules, privacy controls, appeal procedures and robust manual-review workflows.", "", "## 10. Completion Checklist", "", "- [x] Raw source is preserved and its inadequacy is documented.\n- [x] Reproducible synthetic Indian-context fallback is generated and labelled.\n- [x] Stratified 70/15/15 split and train-only transforms are saved.\n- [x] Classical model benchmark, threshold optimisation, INR proxy cost and computational metrics are reported.\n- [x] EDA, ROC/PR, comparison, confusion, feature-importance, fairness and calibration assets are generated.\n- [x] Documentation, presentation and IEEE-style report are generated from the run.\n- [ ] Independent real-insurer validation and governance approval (outside this academic synthetic study).", "", "## References", ""] + REFERENCES
    doc_path = docs / "approach_1_traditional_ml_documentation.md"
    doc_path.write_text("\n".join(lines) + "\n", encoding="utf-8")
    # Evaluation report is a distinct evidence audit; written in one complete operation.
    ev = ["# Approach 1 — Evaluation Evidence and Verification", "", f"Generated: {now}. This report is derived from the benchmark files produced by `src.train`; source data are synthetic educational records.", "", "## Protocol", "", "- Split: stratified 70% train / 15% validation / 15% held-out test.\n- Hyperparameter selection: five-fold stratified CV on training data, mean F2.\n- Threshold selection: validation-only maximum F2 scan.\n- Final comparison: one untouched test partition.\n- Ranking: F2, then ROC-AUC.\n- Cost: FN equals synthetic claim amount; FP=₹3,500 configured review/friction proxy.\n- Warning: no p-value or synthetic metric establishes real-world insurer efficacy.", "", "## Complete benchmark", ""]
    ev += _md_table(public, [c for c in ["algorithm","accuracy","precision","recall","f1","f2","auc_roc","auc_pr","mcc","brier_score","threshold","training_time_seconds","prediction_time_per_sample_ms","model_size_kb","tuned_hyperparameters","true_negative","false_positive","false_negative","true_positive"] if c in public])
    ev += ["", "## Per-model tuning evidence", ""]
    for _, row in public.iterrows():
        ev += [f"### {row['algorithm']}", f"- CV F2: {_fmt(row['cv_f2_mean'])} ± {_fmt(row['cv_f2_std'])}; fold values: `{row['cv_f2_scores']}`.", f"- Validation-selected threshold: {_fmt(row['threshold'])}; validation F2: {_fmt(row['validation_f2'])}.", f"- Best hyperparameters: `{row['best_hyperparameters']}`.", f"- Training/tuning: {_fmt(row['training_time_seconds'],2)} seconds; held-out latency: {_fmt(row['prediction_time_per_sample_ms'],4)} ms/claim; serialized estimator: {_fmt(row['model_size_kb'],1)} KB.", f"- Held-out confusion matrix: TN={row['true_negative']}, FP={row['false_positive']}, FN={row['false_negative']}, TP={row['true_positive']}.", ""]
    ev += ["## Imbalance strategy comparison", ""] + _md_table(imbalance, list(imbalance.columns)) + ["", "## Pairwise significance", ""]
    ev += _md_table(significance, list(significance.columns)) if not significance.empty else ["No paired comparison was applicable."]
    ev += ["", "## Best-model fairness audit", ""] + _md_table(fairness, ["dimension","group","n","accuracy","fpr","fnr","precision","recall","selection_rate","small_group_warning"]) + ["", "## Verification assertions", "", "1. All metrics in this document came from arrays generated by saved models on the test split.\n2. The threshold is selected before test scoring.\n3. Target encoders, imputers, scalers and feature selection are fitted to train rows only.\n4. `benchmark_results.csv`, `fairness_by_group.csv`, `imbalance_strategy_comparison.csv`, `significance_tests.csv` and model metadata are the machine-readable source of truth.\n5. Presentation and PDF values are derived from the same benchmark at generation time.", "", "## Interpretation boundary", "", "Results are valid only for this deterministic synthetic simulation. They cannot be used to approve, reject or price a real medical insurance claim. A real deployment requires external validation, insurer/legal policy mapping, data-protection controls, model-risk governance, calibrated monitoring and accountable human review."]
    ev_path = evaluation / "approach_1_evaluation_report.md"
    ev_path.write_text("\n".join(ev) + "\n", encoding="utf-8")
    return doc_path, ev_path


def _add_title(slide: Any, title: str, subtitle: str | None = None) -> None:
    """Add a consistent presentation heading."""
    box = slide.shapes.add_textbox(Inches(.55), Inches(.32), Inches(12.2), Inches(.55))
    text = box.text_frame; p = text.paragraphs[0]; p.text = title; p.font.size = Pt(25); p.font.bold = True; p.font.color.rgb = RGBColor(23,50,77)
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(.58), Inches(.90), Inches(11.7), Inches(.3)); p = sub.text_frame.paragraphs[0]; p.text = subtitle; p.font.size = Pt(10); p.font.color.rgb = RGBColor(90,100,110)


def _add_footer(slide: Any, index: int) -> None:
    """Add IIIT Dharwad/footer identity and slide number."""
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(.45), Inches(7.18), Inches(12.4), Inches(.05)); line.fill.solid(); line.fill.fore_color.rgb=RGBColor(199,61,79); line.line.fill.background()
    box=slide.shapes.add_textbox(Inches(.55), Inches(7.23), Inches(11.8), Inches(.18)); p=box.text_frame.paragraphs[0]; p.text=f"IIIT Dharwad | Medical Insurance Claim Fraud Detection | Approach 1"; p.font.size=Pt(8); p.font.color.rgb=RGBColor(90,100,110)
    number=slide.shapes.add_textbox(Inches(12.1), Inches(7.23), Inches(.35), Inches(.18)); p=number.text_frame.paragraphs[0];p.text=str(index);p.font.size=Pt(8);p.alignment=PP_ALIGN.RIGHT


def _add_bullets(slide: Any, items: list[str], x: float=.85, y: float=1.35, w: float=5.8, h: float=5.3, size: int=17) -> None:
    """Render concise, readable bullet text."""
    box=slide.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=box.text_frame; tf.word_wrap=True
    for i,item in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.text=item; p.level=0; p.font.size=Pt(size); p.font.color.rgb=RGBColor(35,48,61); p.space_after=Pt(10)


def _add_image(slide: Any, path: Path, x: float, y: float, w: float, h: float) -> None:
    """Add an image if it exists, otherwise render an explanatory placeholder."""
    if path.exists():
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    else:
        shape=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h));shape.fill.solid();shape.fill.fore_color.rgb=RGBColor(238,242,246);shape.line.color.rgb=RGBColor(160,170,180)
        p=shape.text_frame.paragraphs[0];p.text=f"Asset unavailable\n{path.name}";p.alignment=PP_ALIGN.CENTER;p.font.size=Pt(14)


def generate_presentation(config: dict[str, Any], benchmark: pd.DataFrame, quality: dict[str, Any], source_audit: dict[str, Any]) -> Path:
    """Generate a professional 20-slide Approach-1 presentation from actual outputs."""
    public=benchmark.drop(columns=[c for c in benchmark if c.startswith("_")],errors="ignore").sort_values(["f2","auc_roc"],ascending=False).reset_index(drop=True); best=public.iloc[0]
    visual=Path(config["paths"]["visualization_dir"]); out=Path(config["paths"]["presentation_dir"]);out.mkdir(parents=True,exist_ok=True)
    prs=Presentation();prs.slide_width=Inches(13.333);prs.slide_height=Inches(7.5); blank=prs.slide_layouts[6]
    def slide(title: str, subtitle: str | None=None):
        s=prs.slides.add_slide(blank); bg=s.background.fill;bg.solid();bg.fore_color.rgb=RGBColor(250,252,254);_add_title(s,title,subtitle);_add_footer(s,len(prs.slides));return s
    s=slide("Medical Insurance Claim Fraud Detection","Approach 1: Explainable Traditional Machine Learning Baseline")
    banner=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(.55),Inches(1.22),Inches(12.2),Inches(1.15));banner.fill.solid();banner.fill.fore_color.rgb=RGBColor(23,50,77);banner.line.fill.background();p=banner.text_frame.paragraphs[0];p.text="AI-Driven Claim Verification & Responsible Fraud Screening";p.font.size=Pt(27);p.font.bold=True;p.font.color.rgb=RGBColor(255,255,255);p.alignment=PP_ALIGN.CENTER
    _add_bullets(s,["IIIT Dharwad — Department of Data Science and AI","Faculty Adviser: Ramesh Athe","B Varshith (23BDS011) | M Jagadeshwar (23BDS033) | J Ganesh (23BDS024)","Academic mini-project | 7 August 2026"],x=1.3,y=3.05,w=10.8,h=2.8,size=18)
    s=slide("Problem Statement","Fraud screening must protect genuine policyholders without automated adverse decisions");_add_bullets(s,["Medical claims can involve inflated bills, duplicate episodes, timing manipulation and inconsistent documents.","Indian context matters: family floater/group covers, Ayushman Bharat/ECHS, allopathic and Ayurvedic care, GST and regional cost variation.","Goal: classify historic tabular claims as high-risk or legitimate, then route high-risk cases to human evidence review.","A score is not a denial: policy eligibility and due process remain human responsibilities."],w=6.0);_add_image(s,visual/"eda"/"amount_vs_duration.png",7.0,1.35,5.5,4.8)
    s=slide("Objectives and Scope","What this baseline does—and does not—claim");_add_bullets(s,["Benchmark diverse classical models with recall-weighted F2 selection.","Build leakage-aware preprocessing and interpretable feature evidence.","Audit subgroup error rates, calibration, INR-oriented screening costs and computational cost.","Use a transparent synthetic fallback because the bundled file is not adequate Indian claim data.","Not a production settlement engine or estimate of real insurer fraud performance."],w=11.5,size=18)
    s=slide("Dataset Decision","Preserve the supplied workbook; do not misrepresent it");_add_bullets(s,[f"Supplied workbook: {source_audit['rows']:,} records (below the ≥10,000 adequacy threshold).", "Adequacy review: generic locations and missing Indian policy/history fields.","Fallback: deterministic Indian-context synthetic claims in INR, labelled for educational benchmarking.","Raw reference remains untouched with checksum and metadata."],w=6.3);_add_image(s,visual/"eda"/"class_and_claim_type_mix.png",7.1,1.45,5.1,4.4)
    s=slide("Synthetic Population and Split","Same frozen split supports fair later ML-vs-DL comparison");_add_bullets(s,[f"After exact duplicate removal: {quality['rows_after_cleaning']:,} claims.",f"Stratified train/validation/test: {quality['stratified_split']['train_rows']:,} / {quality['stratified_split']['validation_rows']:,} / {quality['stratified_split']['test_rows']:,}.","Features: demographic context, policy terms, claim/provider details, temporal history and regional treatment baseline.","Fraud rate is synthetic; values cannot be generalized to an insurer."],w=6.0);_add_image(s,visual/"eda"/"claim_amount_distribution.png",6.7,1.35,5.7,4.5)
    s=slide("Exploratory Data Analysis","Right-skewed INR costs and regional/treatment heterogeneity");_add_image(s,visual/"eda"/"numeric_correlation_heatmap.png",.75,1.35,6.0,4.85);_add_image(s,visual/"eda"/"fraud_rate_state_treatment_heatmap.png",6.95,1.35,5.65,4.85)
    s=slide("Leakage-Aware Preprocessing","All learned transforms fit on training data only");_add_image(s,visual/"technical"/"traditional_ml_pipeline_diagram.png",.85,1.35,11.6,2.1);_add_bullets(s,["Exact duplicate removal; missingness audit; plausible outliers retained as potential signal.","Median/mode imputation, robust scaling, one-hot low-cardinality encoding and smoothed target encoding for high-cardinality features.",f"Mutual-information selection: {quality['feature_selection']['features_before_selection']} transformed inputs → {quality['feature_selection']['features_after_selection']} selected inputs.","Validation chooses threshold; held-out test is evaluated once."],y=3.85,w=11.6,size=16)
    s=slide("Feature Engineering","Domain knowledge expressed as auditable features");_add_bullets(s,["Claim-to-premium and claim-to-sum-insured ratios", "Regional treatment-cost deviation and amount per hospital day", "Days relative to waiting-period end and claim-frequency intensity", "Current-to-history amount ratio; provider-risk × distance interaction", "Limited squared amount/stay terms—no uncontrolled feature explosion"],w=6.0);_add_image(s,visual/"eda"/"correlation_network.png",7.0,1.35,5.3,4.6)
    s=slide("Model Suite","Thirteen+ classical classifiers under one F2-first protocol");_add_bullets(s,["Linear: L1/L2 Logistic Regression; QDA; Gaussian and Multinomial Naive Bayes", "Trees/ensembles: Decision Tree, Random Forest, Histogram GB, AdaBoost, XGBoost, LightGBM", "Geometry: Linear/RBF SVM and KNN", "Bridge model: shallow neural network with Adam/early stopping", "5-fold stratified CV; grid search for small spaces and random search for broader tree/NN spaces"],w=11.5,size=17)
    s=slide("Class-Imbalance Strategy","Recall is important, but must be balanced against false flags");_add_bullets(s,["Compared: class-weighted learning, random undersampling, Tomek links, SMOTE and SMOTEENN.","Sampling is fit only on train data. Validation determines F2 threshold.","The full benchmark uses consistent class-aware settings; results disclose precision, recall, PR-AUC, MCC and INR proxy cost."],w=6.0);# simple table rendered as bullets
    for i,(_,r) in enumerate(pd.read_csv(Path(config['paths']['evaluation_dir'])/'imbalance_strategy_comparison.csv').head(5).iterrows()):
        box=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(7),Inches(1.45+i*.75),Inches(4.9),Inches(.57));box.fill.solid();box.fill.fore_color.rgb=RGBColor(235,242,247);box.line.color.rgb=RGBColor(180,195,205);p=box.text_frame.paragraphs[0];p.text=f"{r['strategy']}: validation F2={r['validation_f2']:.3f}, recall={r['validation_recall']:.3f}";p.font.size=Pt(12)
    s=slide("Evaluation Design","Why F2, PR-AUC, costs and paired tests are all necessary");_add_bullets(s,["F2 is primary: false negatives receive greater weight because missed suspicious claims can be costly.","PR-AUC complements ROC-AUC when fraud is minority class; MCC remains informative under imbalance.","Threshold scans run on validation only—not default 0.50.","McNemar exact test compares same-test predictions; Wilcoxon compares CV folds.","Costs are illustrative: missed synthetic fraud amount + ₹3,500 review/friction per false flag."],w=11.5,size=17)
    s=slide("Held-out Benchmark","Ranked by F2, then ROC-AUC; test set untouched until final evaluation");_add_image(s,visual/"model_comparison"/"grouped_metric_comparison.png",.6,1.25,12.0,5.3)
    s=slide("ROC and PR Curves","PR curves reveal positive-class performance under imbalance");_add_image(s,visual/"model_comparison"/"roc_curves_all_models.png",.55,1.2,6.15,5.3);_add_image(s,visual/"model_comparison"/"precision_recall_curves_all_models.png",6.75,1.2,6.05,5.3)
    s=slide("Best-Performer Summary",f"Held-out leader: {best['algorithm']}");_add_bullets(s,[f"F2 = {best['f2']:.3f}; Recall = {best['recall']:.3f}; Precision = {best['precision']:.3f}",f"ROC-AUC = {best['auc_roc']:.3f}; PR-AUC = {best['auc_pr']:.3f}; MCC = {best['mcc']:.3f}",f"Validation-selected threshold = {best['threshold']:.3f}",f"Training/tuning = {best['training_time_seconds']:.2f}s; latency = {best['prediction_time_per_sample_ms']:.4f} ms/claim",f"This result describes the synthetic test set only."],w=6.1,size=18);_add_image(s,visual/"model_comparison"/"top_five_radar.png",7.0,1.3,5.3,4.9)
    s=slide("Computational Efficiency","Operational choice needs more than a single score");_add_image(s,visual/"model_comparison"/"training_time_vs_accuracy.png",.75,1.35,6.3,4.7);_add_image(s,visual/"model_comparison"/"confusion_matrices_all_models.png",7.25,1.35,5.25,4.7)
    s=slide("Model Interpretability","Associations guide review; they do not prove fraud");importance_assets=sorted((visual/"interpretability").glob("feature_importance*.png"));_add_image(s,importance_assets[0] if importance_assets else visual/"technical"/"traditional_ml_pipeline_diagram.png",.65,1.25,6.35,4.95);_add_bullets(s,["Coefficients/rules/tree importances identify inputs associated with synthetic risk.","A claimant-facing explanation must cite verifiable documents/policy clauses, not merely a model feature.","High-risk recommendation → human investigator verifies billing, coverage and clinical evidence.","No demographic attribute should be an adverse-decision reason."],x=7.3,y=1.5,w=5.1,h=4.9,size=16)
    s=slide("Fairness and Calibration","Audit outcomes across groups; investigate disparities");_add_image(s,visual/"fairness"/"fairness_fnr_groups.png",.8,1.35,6.1,4.7);_add_image(s,visual/"fairness"/"calibration_reliability_diagram.png",7.1,1.35,5.3,4.7)
    s=slide("Strengths, Limitations and Approach Comparison","Responsible interpretation of an academic baseline");_add_bullets(s,["Strengths: reproducibility, Indian-context assumptions, broad model suite, threshold/cost/fairness evidence and explainability assets.","Limits: labels and mechanisms are synthetic; random claim split can retain correlated entities; provider history is a proxy.","Approach 2 will test learned embeddings/temporal neural representations on this frozen split; Approach 3 will add document/RAG evidence and human checkpoints.","No model output may automatically reject a person’s claim; production needs governed data, validation, privacy/security, monitoring, appeal and human accountability."],w=11.5,size=16)
    s=slide("Conclusion, Future Work and References","Key takeaways and reproducibility direction");_add_bullets(s,[f"{len(public)} models were evaluated under a consistent, F2-led synthetic protocol; the leader is {best['algorithm']} (F2 {best['f2']:.3f}).", "Validation-selected thresholds and INR proxy cost expose operational trade-offs; fairness/calibration/interpretability are required alongside discrimination metrics.","Next: entity-disjoint and temporal validation, insurer-approved labels, policy-specific coverage rules, Approach-2 neural tests and Approach-3 evidence-grounded review.","References include Phua et al. (2010), Joudaki et al. (2015), Chawla et al. (2002), Chen & Guestrin (2016), Ke et al. (2017), Lundberg & Lee (2017), Hardt et al. (2016) and IRDAI (2024). Full list in documentation."],w=11.5,size=15)
    s=slide("Thank You","Questions and discussion");_add_bullets(s,["Medical Insurance Claim Fraud Detection", "IIIT Dharwad — Department of Data Science and AI", "Faculty Adviser: Ramesh Athe", "Team: B Varshith | M Jagadeshwar | J Ganesh", "Contact: project team via IIIT Dharwad academic channel"],x=1.5,y=2,w=10.5,h=3.8,size=22)
    path=out/"approach_1_traditional_ml_presentation.pptx";prs.save(path);return path


def generate_ieee_pdf(config: dict[str, Any], benchmark: pd.DataFrame, source_audit: dict[str, Any], visual_assets: list[Path]) -> Path:
    """Create a substantive IEEE-inspired two-column PDF research report with figures/tables.

    Args:
        config: Loaded configuration.
        benchmark: Current benchmark table.
        source_audit: Raw-source audit.
        visual_assets: Created plot paths to reuse as figures.

    Returns:
        PDF report path.
    """
    public=benchmark.drop(columns=[c for c in benchmark if c.startswith("_")],errors="ignore").sort_values(["f2","auc_roc"],ascending=False).reset_index(drop=True);best=public.iloc[0]
    output=Path(config["paths"]["reports_dir"]);output.mkdir(parents=True,exist_ok=True);path=output/"approach_1_ieee_style_research_report.pdf"
    styles=getSampleStyleSheet();styles.add(ParagraphStyle(name="PaperTitle",parent=styles["Title"],fontName="Helvetica-Bold",fontSize=17,leading=20,alignment=TA_CENTER,spaceAfter=7));styles.add(ParagraphStyle(name="Authors",parent=styles["Normal"],fontSize=9,leading=11,alignment=TA_CENTER,spaceAfter=9));styles.add(ParagraphStyle(name="PaperBody",parent=styles["BodyText"],fontSize=8.2,leading=10.4,alignment=TA_JUSTIFY,spaceAfter=5));styles.add(ParagraphStyle(name="PaperHead",parent=styles["Heading2"],fontName="Helvetica-Bold",fontSize=10,leading=12,spaceBefore=7,spaceAfter=4));styles.add(ParagraphStyle(name="Caption",parent=styles["Normal"],fontSize=7.2,leading=8.5,alignment=TA_CENTER,spaceAfter=6))
    class Doc(BaseDocTemplate):
        pass
    doc=Doc(str(path),pagesize=A4,rightMargin=1.25*cm,leftMargin=1.25*cm,topMargin=1.1*cm,bottomMargin=1.1*cm)
    width,height=A4;gap=.45*cm;column=(width-doc.leftMargin-doc.rightMargin-gap)/2
    frame1=Frame(doc.leftMargin,doc.bottomMargin,column,height-doc.topMargin-doc.bottomMargin,id="col1");frame2=Frame(doc.leftMargin+column+gap,doc.bottomMargin,column,height-doc.topMargin-doc.bottomMargin,id="col2")
    def decorate(canvas, doc_obj):
        canvas.saveState();canvas.setFont("Helvetica",6.5);canvas.setFillColor(colors.HexColor("#555555"));canvas.drawString(doc.leftMargin, .65*cm,"IIIT Dharwad | Medical Insurance Claim Fraud Detection | Traditional ML");canvas.drawRightString(width-doc.rightMargin,.65*cm,f"Page {doc_obj.page}");canvas.restoreState()
    from reportlab.platypus import PageTemplate
    doc.addPageTemplates([PageTemplate(id="two",frames=[frame1,frame2],onPage=decorate)])
    story=[Paragraph("Medical Insurance Claim Fraud Detection: A Reproducible Traditional Machine Learning Baseline for an Indian-Context Synthetic Study",styles["PaperTitle"]),Paragraph("B Varshith (23BDS011), M Jagadeshwar (23BDS033), J Ganesh (23BDS024)<br/>Department of Data Science and AI, IIIT Dharwad, India<br/>Faculty Adviser: Ramesh Athe",styles["Authors"])]
    abstract=(f"<b>Abstract—</b>This paper presents an interpretable traditional-machine-learning baseline for medical insurance claim fraud screening in an Indian context. The bundled workbook was preserved, audited and not used as the primary corpus because it contains only {source_audit['rows']:,} records and lacks Indian policy, geography and historical-claim fields. A deterministic, explicitly synthetic INR-denominated claim population was generated with policy structures, regional/tier cost variation, allopathic and Ayurvedic treatment types, and probabilistic fraud labels. Train-only preprocessing, five-fold F2-led tuning, validation-only threshold selection and a held-out test comparison were applied to {len(public)} classifiers. The best held-out synthetic result was {best['algorithm']} with F2={best['f2']:.3f}, recall={best['recall']:.3f} and PR-AUC={best['auc_pr']:.3f}. Results are not estimates of real insurer effectiveness. The study contributes a transparent software and reporting baseline that includes class-imbalance comparison, INR-oriented screening-cost assumptions, calibration, fairness audit and interpretable feature assets. <b>Keywords—</b>health insurance fraud, imbalanced learning, India, explainable AI, F2 score, synthetic data.")
    story += [Paragraph(abstract,styles["PaperBody"]),Paragraph("I. INTRODUCTION",styles["PaperHead"]),Paragraph("Medical insurance screening must distinguish suspicious patterns from legitimate, costly care while preserving fairness and due process. Indian products include individual, family-floater and group covers, government schemes, tiered provider costs and diverse treatment practices. A classifier therefore provides only one risk signal; a human reviewer must validate documents, medical necessity and policy clauses. This project establishes the transparent tabular baseline that later deep-learning and evidence-grounded agentic approaches can compare against.",styles["PaperBody"]),Paragraph("II. RELATED WORK",styles["PaperHead"]),Paragraph("Fraud detection literature stresses class imbalance, evolving behavioural patterns and careful evaluation [1]–[4]. Health-care studies use provider and claim features but face label and policy-context limitations [5], [6]. SMOTE and cost-sensitive learning address minority detection [7]–[9], while forests and boosted trees remain strong structured-data baselines [10]–[13]. Explainability and fairness work warns that predictive features do not constitute causal or legal evidence [17], [20], [21].",styles["PaperBody"]),Paragraph("III. DATA AND METHODS",styles["PaperHead"]),Paragraph("The raw workbook is archived with a checksum and a documented adequacy decision. The fallback generator creates right-skewed INR amounts, 5–15% target fraud prevalence, state/city/tier baselines, policy duration/waiting periods/co-pay, claim type, historical frequencies and provider proxies. Exact duplicates are detected and removed. The 70/15/15 stratified split is frozen. Imputation, robust scaling, one-hot encoding, smoothed target encoding and mutual-information selection are fitted on the training partition only.",styles["PaperBody"])]
    table_data=[["Metric","Protocol"],["Primary selection","5-fold train CV F2"],["Threshold","Validation-only maximum F2"],["Final assessment","Untouched stratified test set"],["Cost proxy","FN=synthetic claim INR; FP=₹3,500"],["Fairness","FPR/FNR/accuracy by audit groups"]]
    t=Table(table_data,colWidths=[column*.38,column*.57]);t.setStyle(TableStyle([("BACKGROUND",(0,0),(-1,0),colors.HexColor("#17324d")),("TEXTCOLOR",(0,0),(-1,0),colors.white),("FONTNAME",(0,0),(-1,0),"Helvetica-Bold"),("FONTSIZE",(0,0),(-1,-1),6.8),("GRID",(0,0),(-1,-1),.25,colors.grey),("VALIGN",(0,0),(-1,-1),"TOP"),("ROWBACKGROUNDS",(0,1),(-1,-1),[colors.white,colors.HexColor("#f3f6f8")])]))
    story += [Paragraph("Table I. Experimental protocol",styles["Caption"]),t,Spacer(1,6),Paragraph("IV. RESULTS AND ANALYSIS",styles["PaperHead"])]
    bench_data=[["Model","F2","Rec.","PR-AUC","AUC","ms"]]+[[str(r["algorithm"])[:22],f"{r['f2']:.3f}",f"{r['recall']:.3f}",f"{r['auc_pr']:.3f}",f"{r['auc_roc']:.3f}",f"{r['prediction_time_per_sample_ms']:.2f}"] for _,r in public.iterrows()]
    tb=Table(bench_data,colWidths=[column*.31,column*.12,column*.12,column*.15,column*.13,column*.12]);tb.setStyle(TableStyle([("BACKGROUND",(0,0),(-1,0),colors.HexColor("#17324d")),("TEXTCOLOR",(0,0),(-1,0),colors.white),("FONTNAME",(0,0),(-1,0),"Helvetica-Bold"),("FONTSIZE",(0,0),(-1,-1),5.5),("GRID",(0,0),(-1,-1),.2,colors.grey),("VALIGN",(0,0),(-1,-1),"MIDDLE"),("ROWBACKGROUNDS",(0,1),(-1,-1),[colors.white,colors.HexColor("#f3f6f8")])]))
    story += [Paragraph("Table II. Held-out benchmark (ranked by F2)",styles["Caption"]),tb,Spacer(1,6)]
    adequacy_data=[["Source criterion","Observed status"],["Rows",f"{source_audit['rows']:,}; below the ≥10,000 study criterion"],["Label","Binary ClaimLegitimacy field is present"],["Indian context","Generic locations/identifiers; not adequate"],["Policy/history fields","Required waiting/co-pay/history fields absent"],["Decision","Preserve source; use labelled synthetic fallback"]]
    adequacy_table=Table(adequacy_data,colWidths=[column*.40,column*.55]);adequacy_table.setStyle(TableStyle([("BACKGROUND",(0,0),(-1,0),colors.HexColor("#17324d")),("TEXTCOLOR",(0,0),(-1,0),colors.white),("FONTNAME",(0,0),(-1,0),"Helvetica-Bold"),("FONTSIZE",(0,0),(-1,-1),6.5),("GRID",(0,0),(-1,-1),.2,colors.grey),("VALIGN",(0,0),(-1,-1),"TOP"),("ROWBACKGROUNDS",(0,1),(-1,-1),[colors.white,colors.HexColor("#f3f6f8")])]))
    feature_data=[["Feature family","Examples / controls"],["Policy", "sum insured, premium, co-pay, waiting period"],["Claim", "INR amount, type, treatment, duration, procedure count"],["History", "frequency, time since prior claim, amount baseline"],["Provider/context", "tier, network status, distance, regional baseline"],["Engineered", "cost deviation, utilisation, amount/day; train-only encoding"]]
    feature_table=Table(feature_data,colWidths=[column*.35,column*.60]);feature_table.setStyle(TableStyle([("BACKGROUND",(0,0),(-1,0),colors.HexColor("#17324d")),("TEXTCOLOR",(0,0),(-1,0),colors.white),("FONTNAME",(0,0),(-1,0),"Helvetica-Bold"),("FONTSIZE",(0,0),(-1,-1),6.5),("GRID",(0,0),(-1,-1),.2,colors.grey),("VALIGN",(0,0),(-1,-1),"TOP"),("ROWBACKGROUNDS",(0,1),(-1,-1),[colors.white,colors.HexColor("#f3f6f8")])]))
    cost_data=[["Error outcome","Illustrative handling"],["False negative", "Synthetic fraudulent claim recommended approve; cost equals its INR amount"],["False positive", "Legitimate claim routed to review; ₹3,500 friction/review proxy"],["Threshold", "Selected on validation maximum F2, never using test labels"],["Decision authority", "Human reviewer verifies documents, policy and medical evidence"]]
    cost_table=Table(cost_data,colWidths=[column*.32,column*.63]);cost_table.setStyle(TableStyle([("BACKGROUND",(0,0),(-1,0),colors.HexColor("#17324d")),("TEXTCOLOR",(0,0),(-1,0),colors.white),("FONTNAME",(0,0),(-1,0),"Helvetica-Bold"),("FONTSIZE",(0,0),(-1,-1),6.5),("GRID",(0,0),(-1,-1),.2,colors.grey),("VALIGN",(0,0),(-1,-1),"TOP"),("ROWBACKGROUNDS",(0,1),(-1,-1),[colors.white,colors.HexColor("#f3f6f8")])]))
    story += [Paragraph("Table III. Source adequacy decision",styles["Caption"]),adequacy_table,Spacer(1,5),Paragraph("Table IV. Feature families and leakage controls",styles["Caption"]),feature_table,Spacer(1,5),Paragraph("Table V. Illustrative operational cost assumptions",styles["Caption"]),cost_table,Spacer(1,6),Paragraph(f"The F2-leading synthetic model is {best['algorithm']}. Its threshold of {best['threshold']:.3f} was selected from validation probabilities, avoiding a default 0.50 assumption. Precision, recall, PR-AUC, MCC, calibration and INR proxy cost must be read together; a higher F2 alone does not justify automatic rejection.",styles["PaperBody"])]
    # at least five figure references, subject to existing asset availability.
    figure_candidates=[p for p in visual_assets if p.exists()][:6]
    captions=["Fig. 1. Synthetic claim amount distribution by class.","Fig. 2. Numeric correlation structure.","Fig. 3. Held-out all-model metric comparison.","Fig. 4. Held-out ROC comparison.","Fig. 5. Held-out precision–recall comparison.","Fig. 6. Fairness/calibration audit."]
    for idx,image_path in enumerate(figure_candidates):
        story += [KeepTogether([Image(str(image_path),width=column*.93,height=column*.59,kind="proportional"),Paragraph(captions[idx],styles["Caption"])])]
    story += [Paragraph("V. FAIRNESS, EXPLAINABILITY AND DISCUSSION",styles["PaperHead"]),Paragraph("The held-out audit reports group accuracy, FPR, FNR, precision, recall and selection rate across gender, age, geography, income and treatment. Small groups are marked rather than overinterpreted. A material disparity triggers data-quality review and a mitigation comparison; protected attributes must never act as unreviewed grounds for a denial. Native feature importances and coefficients are associations within the synthetic generator. They must be translated into reviewable evidence, such as a bill discrepancy or policy clause, before communicating with a policyholder.",styles["PaperBody"]),Paragraph("VI. LIMITATIONS AND FUTURE WORK",styles["PaperHead"]),Paragraph("Synthetic labels and prices prevent operational claims. Future work requires governed de-identified insurer data, temporal/entity-disjoint validation, data drift monitoring, external calibration, policy-specific cost matrices, documented appeals and human-review performance measurement. Deep-learning models should use the same frozen test protocol; the agentic workflow should add document/RAG evidence without representing LLM reasoning as factual proof.",styles["PaperBody"]),Paragraph("VII. CONCLUSION",styles["PaperHead"]),Paragraph("This work provides a reproducible and transparent classical baseline for educational Indian-context claim fraud screening. It makes data limitations visible, prioritizes recall-aware evidence, and packages performance with cost, interpretability and fairness audits. It is a foundation for accountable research, not an automatic insurance decision system.",styles["PaperBody"]),Paragraph("ACKNOWLEDGMENTS",styles["PaperHead"]),Paragraph("The authors thank Faculty Adviser Ramesh Athe for guidance and the Department of Data Science and AI, IIIT Dharwad, for academic support.",styles["PaperBody"]),Paragraph("REFERENCES",styles["PaperHead"])]
    story += [Paragraph(ref,styles["PaperBody"]) for ref in REFERENCES]
    doc.build(story);return path
