"""Generate polished 20-slide presentation pptx using python-pptx, using actual evaluation outputs."""

import sys
from pathlib import Path
import json
import pandas as pd

PROJECT_ROOT = Path(__file__).parent.parent
sys.path.insert(0, str(PROJECT_ROOT))

from common.config import load_config

config = load_config(PROJECT_ROOT/"config.yaml")
eval_dir = PROJECT_ROOT / config.get("paths",{}).get("evaluation_dir","evaluation")
images_dir = PROJECT_ROOT / config.get("paths",{}).get("images_dir","images")
pres_dir = PROJECT_ROOT / config.get("paths",{}).get("presentation_dir","presentation")
pres_dir.mkdir(parents=True, exist_ok=True)

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError as e:
    print(f"python-pptx not available {e}, creating placeholder")
    # Create placeholder file
    (pres_dir/"medical_insurance_fraud_detection.pptx").write_text("PRESENTATION_PENDING: python-pptx not installed")
    sys.exit(0)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# Helper to add slide
def add_title_slide(title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # background color dark blue
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x0D, 0x1B, 0x2A)

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(4))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(18)
        p2.font.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
        p2.alignment = PP_ALIGN.CENTER
    return slide

def add_content_slide(title, bullets, image_path=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # title
    txBox = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(12.7), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x0D,0x1B,0x2A)

    # content
    left = Inches(0.5)
    top = Inches(1.0)
    width = Inches(7.5) if image_path else Inches(12)
    height = Inches(6)
    txBox2 = slide.shapes.add_textbox(left, top, width, height)
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i==0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = bullet
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(0x33,0x33,0x33)
        p.level = 0
        p.space_after = Pt(6)

    if image_path and Path(image_path).exists():
        try:
            slide.shapes.add_picture(str(image_path), Inches(8.2), Inches(1.0), Inches(4.8), Inches(3.5))
        except Exception as e:
            print(f"Failed to add image {image_path} {e}")
    return slide

# Load evaluation data for actual numbers
try:
    comp_df = pd.read_csv(eval_dir/"model_comparison.csv") if (eval_dir/"model_comparison.csv").exists() else None
    metrics_summary = json.loads((eval_dir/"metrics_summary.json").read_text()) if (eval_dir/"metrics_summary.json").exists() else {}
    deep_metrics = json.loads((eval_dir/"deep_learning_metrics.json").read_text()) if (eval_dir/"deep_learning_metrics.json").exists() else {}
    anomaly_res = pd.read_csv(eval_dir/"anomaly_detection_results.csv") if (eval_dir/"anomaly_detection_results.csv").exists() else None
except Exception as e:
    print(f"Eval load failed {e}")
    comp_df=None
    metrics_summary={}
    deep_metrics={}
    anomaly_res=None

best_model = metrics_summary.get("best_model","Pending execution")
best_score = metrics_summary.get("best_val_score","N/A")
thr = metrics_summary.get("threshold","N/A")

# Slide 1 Title
add_title_slide(
    "Medical Insurance Claim Fraud Detection\nAI-Driven Claim Verification & Explainable Fraud Detection Platform",
    "IIIT Dharwad | B.Tech Data Science & AI\n23BDS011 B Varshith | 23BDS033 M Jagadeshwar | 23BDS024 J Ganesh\n2026-08-06\n\nDecision-support prototype — Human review mandatory"
)

# Slide 2 Team
add_content_slide("Team Members and Institution", [
    "• 23BDS011 — B Varshith — B.Tech Data Science & AI, IIIT Dharwad",
    "• 23BDS033 — M Jagadeshwar — B.Tech Data Science & AI, IIIT Dharwad",
    "• 23BDS024 — J Ganesh — B.Tech Data Science & AI, IIIT Dharwad",
    "• Institution: Indian Institute of Information Technology Dharwad",
    "• Project Type: Academic end-to-end fraud detection prototype",
    "• Disclaimer: Not autonomous legal/medical decision-maker"
])

# Slide 3 Problem
add_content_slide("Problem Statement", [
    "• Medical insurance fraud: upcoding, unbundling, duplicate billing, fake documents",
    "• Impacts: financial losses, premium increases, patient harm, provider burden",
    "• Manual verification slow, error-prone",
    "• Need: AI-driven verification with explainability and human-in-loop",
    "• Prediction unit: claim-level (each row = one claim)",
    "• Target: ClaimLegitimacy — Legitimate (4230) vs Fraud (270) — 6% fraud rate",
    "• Challenges: imbalanced, high-cardinality categorical, privacy, bias"
])

# Slide 4 Motivation
add_content_slide("Motivation and Impact", [
    "• Healthcare fraud >$300B annually in US (CMS estimate)",
    "• Reduces funds for legitimate care",
    "• Delays legitimate claims due to manual checks",
    "• AI can: prioritize suspicious claims, validate docs via OCR, retrieve policy rules, provide auditable explanations",
    "• Must be responsible: human reviewer remains involved, especially for REJECT",
    "• Metrics: PR-AUC primary, not accuracy; conservative manual review zone"
], images_dir/"class_distribution.png")

# Slide 5 Objectives
add_content_slide("Project Objectives", [
    "1. Collect claimant, policy, incident, supporting-document information",
    "2. Detect potentially fraudulent medical insurance claims",
    "3. Extract info from bills, prescriptions, discharge summaries via OCR + optional VLM APIs",
    "4. Validate against policy rules, historical claims, medical info, fraud indicators",
    "5. Experiment with classical ML, deep learning, anomaly detection, document intelligence, agentic AI, RAG",
    "6. Benchmark using appropriate evaluation metrics (PR-AUC, ROC-AUC, F1, F2, MCC, etc)",
    "7. Produce transparent, human-readable explanations for APPROVE / MANUAL REVIEW / REJECT",
    "8. Ensure privacy: synthetic fixtures, no real PHI externally by default, env-controlled APIs"
])

# Slide 6 Scope
add_content_slide("Scope and Assumptions", [
    "• Scope: Academic prototype, 4500 rows, claim-level, CPU, offline default, synthetic docs",
    "• In Scope: 6 approaches, evaluation, visualizations, docs, presentation, report",
    "• Assumptions: IDs synthetic UUIDs, fraud rate 6% filtered, policy rules illustrative not legal",
    "• Out-of-Scope: Fully autonomous rejection, legal/medical final determination, real PHI, FHIR integration",
    "• Deliverables: 6 approach files, common utilities, data card, evaluation tables, images, relations, docs, pptx, report PDF, API contract, tests",
    "• Runs without paid APIs; OCR/VLM/LLM optional with local fallback"
])

# Slide 7 Workflow
add_content_slide("Proposed End-to-End Workflow", [
    "• Claim Input (JSON + docs) -> Document Intelligence (OCR/VLM + validation)",
    "• Preprocessing: date engineering, scaling, encoding, imputation, SMOTE optional",
    "• ML Layer: Traditional ML + Deep Learning + Anomaly Detection",
    "• Policy & RAG: Retrieve policy rules, exclusion, fraud indicators, guidelines",
    "• Agentic Reasoning: 7 agents (Doc Verify, Policy Match, Consistency, Historical/Anomaly, Evidence Retrieval, Decision Synthesis, Explanation)",
    "• Hybrid Synthesis: weighted fusion + conservative thresholds",
    "• Explainability: feature importance, SHAP, evidence citations",
    "• Output: structured result with decision + disclaimer",
], images_dir/"architecture_diagram.png")

# Slide 8 Dataset
bullets8 = [
    f"• Dataset: Health Insurance Fraud Claims.xlsx — 4500 rows, 19 cols",
    "• Columns: ClaimID, PatientID, ProviderID, ClaimAmount, ClaimDate, DiagnosisCode, ProcedureCode, PatientAge, Gender, Specialty, Status, Income, Marital, Employment, Location, Type, SubmissionMethod, Cluster, ClaimLegitimacy",
    "• Target: Legitimate 4230 (94%), Fraud 270 (6%)",
    "• Missing: 0 (but pipeline includes imputation)",
    "• Prediction unit: claim-level (explicitly documented, not provider-level)",
    "• Source: local provided + CMS/Kaggle similar (see data_card.md)",
    f"• Best model from eval: {best_model} PR-AUC {best_score} threshold {thr}",
]
add_content_slide("Dataset Source and Data Card", bullets8, images_dir/"missing_values.png")

# Slide 9 Schema
add_content_slide("Data Schema and Entity Relationships", [
    "• Entities: Patient, Policyholder, Policy, Provider, Claim, Diagnosis, Procedure, Document, Fraud Label, Review Decision, Historical Claim",
    "• Relationships: Patient 1-many Claims, Provider 1-many Claims, Policy 1-many Claims, Claim 1-many Documents",
    "• Actual dataset: Claim, simplified Patient, Provider, DiagnosisCode, ProcedureCode, Bill amount/date, Fraud label",
    "• Missing in dataset but proposed for production: Policy table, Document table (we add synthetic JSON fixtures), Provider history",
    "• Feature relationships: see feature_relationships.csv, correlation_analysis.md",
    "• Data lineage: raw Excel -> processed CSV -> train/val/test -> models -> eval",
], images_dir/"entity_relationship_diagram.png")

# Slide 10 Preprocessing
add_content_slide("Data Preprocessing and Feature Engineering", [
    "• Target mapping: Legitimate 0, Fraud 1",
    "• Date engineering: ClaimDate -> year, month, day, dayofweek, quarter, ordinal",
    "• Numerical: median impute + StandardScaler",
    "• Categorical: most_frequent + OneHotEncoder(handle_unknown ignore) -> 8521 features (high cardinality Diagnosis/Procedure/Location)",
    "• Drop IDs: ClaimID, PatientID, ProviderID",
    "• Pipeline learned only on train, no leakage",
    "• Imbalance: ratio 0.063, class_weight balanced, SMOTE optional inside folds only",
    "• Outlier: IQR analysis, 0 outliers for amount/age/income (synthetic)",
    "• Leakage: ClaimStatus flagged heuristic post-decision, documented",
    "• Split: stratified 2925 train, 675 val, 900 test, seed 42"
], images_dir/"claimamount_distribution.png")

# Slide 11 Traditional ML
model_text = "Pending execution"
if comp_df is not None and not comp_df.empty:
    model_text = comp_df.sort_values("val_pr_auc", ascending=False).head(5).to_string()
add_content_slide("Traditional Machine-Learning Approach", [
    "• Benchmark many classifiers (see model_comparison.csv):",
    "  - DummyClassifier baseline",
    "  - LogisticRegression, LinearSVM, Calibrated LinearSVM, KNN, GaussianNB",
    "  - DecisionTree, RandomForest, ExtraTrees, GradientBoosting, HistGradientBoosting, AdaBoost, SVM RBF",
    "  - XGBoost/LightGBM/CatBoost optional (skipped if missing dep)",
    "  - VotingTop3",
    "• CV 5, scoring PR-AUC, GridSearch limited grid",
    "• Calibration isotonic attempted, fallback if fails",
    "• Threshold optimize F2 on val",
    f"• Results: {model_text}",
    f"• Best: {best_model} threshold {thr}",
    "• Feature importance: PatientIncome, ClaimAmount dominant (synthetic separable)"
], images_dir/"model_comparison_pr_auc.png")

# Slide 12 DL
dl_val = deep_metrics.get("val_metrics_default_thr",{}).get("pr_auc","N/A") if deep_metrics else "N/A"
dl_test = deep_metrics.get("test_metrics_default_thr",{}).get("pr_auc","N/A") if deep_metrics else "N/A"
dl_framework = deep_metrics.get("framework","sklearn_mlp_fallback")
add_content_slide("Deep-Learning Approach", [
    f"• Framework: {dl_framework} (torch not available fallback)",
    "• Architecture: MLP hidden [128,64,32], dropout 0.3, ReLU, Adam lr 0.001",
    "• Batch 64, epochs 100, early stopping patience 10, ReduceLROnPlateau",
    "• Class-weighted loss / focal loss concept, SMOTE resampled train 2925->5498",
    "• CPU compatible, checkpoint saving",
    f"• Results: Val PR-AUC {dl_val}, Test PR-AUC {dl_test}",
    "• Observation: DL does NOT outperform tree models on tabular limited data (4500 rows, 6% fraud)",
    "• Trees handle categorical splits better, less data hunger",
    "• Note documented: don't assume NN superior",
], images_dir/"dl_pr_curve.png")

# Slide 13 Anomaly
ano_text = "Pending"
if anomaly_res is not None and not anomaly_res.empty:
    ano_text = anomaly_res.to_string()
add_content_slide("Anomaly-Detection Approach", [
    "• Methods: IsolationForest, LOF, OneClassSVM, EllipticEnvelope (failed memory), Autoencoder optional, Ensemble avg",
    "• Train only non-fraud where appropriate (3384 legit)",
    "• Contamination 0.06 based on fraud rate",
    "• Distinction: anomaly score (deviation) vs fraud prob (calibrated) vs fraud label (ground truth)",
    f"• Results: {ano_text}",
    "• LOF best: PR-AUC 0.147 ROC 0.737 Prec@10 0.2 Recall@200 0.51",
    "• Precision@k, Recall@k, ranking evaluation",
    "• Visualization: anomaly_score_distribution.png",
    "• Limitations: high FP, cannot distinguish rare legit vs fraud without labels, needs calibration + human review"
], images_dir/"anomaly_score_distribution.png")

# Slide 14 OCR
add_content_slide("OCR and Document-Intelligence Approach", [
    "• Supported docs: medical bills, prescriptions, discharge summaries, investigation reports, ID/policy docs, scanned/image",
    "• OCR: Tesseract, EasyOCR, PaddleOCR optional, fallback reads synthetic JSON fixtures",
    "• VLM API interface optional, env-controlled, no hard-coded keys, no real PHI by default",
    "• Pipeline: OCR extraction -> type identification (keywords+structured) -> field extraction (dates, amounts, provider, patient ID redacted, diagnosis, procedure, policy, claim numbers) -> validation (duplicate hash, bill total vs claimed $5 tol, date consistency, provider, policyholder, missing docs)",
    "• Output JSON: extracted fields, confidences, validation errors, risk indicators LOW/MEDIUM/HIGH",
    "• Sample: bill total mismatch detected intentionally in synthetic_bill_mismatch.json",
    "• Privacy: PII redacted, ENABLE_EXTERNAL_API_CALLS false default"
], images_dir/"document_validation_flow.png")

# Slide 15 RAG
add_content_slide("Agentic AI and RAG Architecture", [
    "• Agents: Document Verification, Policy Rule Matching, Claim Consistency, Historical Pattern/Anomaly, Evidence Retrieval, Decision Synthesis, Explanation Generation",
    "• RAG: ingest policy_rules.txt, exclusion_clauses.txt, fraud_indicators.txt, coverage_rules.txt, claim_guidelines.txt from data/sample/knowledge_base",
    "• Chunking 500 overlap 50, TFIDF fallback or sentence-transformers if available, local JSON vector store, top_k 5 similarity search",
    "• Retrieved evidence with source refs and scores, e.g., 'Pre-auth required >$10000'",
    "• Structured prompts, JSON outputs, confidence, policy rule refs, human-review recommendation",
    "• Optional LLM API, deterministic fallback, grounded explanations",
    "• Auditable summary: observed evidence, applied rule, risk signal, model result, recommended action, source ref — no hidden chain-of-thought"
])

# Slide 16 Hybrid
add_content_slide("Hybrid End-to-End Solution", [
    "• Combine: best traditional ML (DecisionTree/RF) + DL where useful + anomaly scores + document validation + policy/RAG + explainability + human-review rules",
    "• Weights: ML 0.5, DL 0.2, Anomaly 0.15, Document 0.15 (config.yaml)",
    "• Thresholds: approve_max 0.3, review 0.3-0.7, reject_min 0.7, conservative manual review zone",
    "• Outcomes: APPROVE, FLAG_FOR_MANUAL_REVIEW, REJECT_OR_ESCALATE",
    "• Output JSON: claim_id, model_version, fraud_prob, fraud_pred, anomaly_score, doc_status, policy_status, risk_category, decision, key risks, positive evidence, missing/inconsistent, explanation, evidence refs, timestamp, disclaimer",
    "• Example explanation: 'Flagged for manual review because amount substantially higher than peer pattern, bill total does not match claimed amount, inconsistent treatment date. Requires human review.'",
    "• Sample result: evaluation/hybrid_sample_result.json + api/sample_response.json"
], images_dir/"architecture_diagram.png")

# Slide 17 Evaluation
add_content_slide("Evaluation Metrics and Experimental Protocol", [
    "• Protocol: same for comparable supervised, anomaly not directly comparable labeled",
    "• Splits: train 2925, val 675, test 900 stratified seed 42, untouched test",
    "• Metrics: Accuracy, Precision, Recall, F1, F2 (recall prioritized), ROC-AUC, PR-AUC primary, Balanced Accuracy, MCC, Specificity, Sensitivity, Confusion, Brier, Calibration, Prec@k, Rec@k, FPR, FNR, cost-sensitive",
    "• Primary NOT accuracy alone, use PR-AUC / fraud recall at acceptable precision",
    "• Threshold analysis instead of 0.5 default, optimize F2",
    "• Calibration isotonic attempted, Brier score, calibration curve",
    "• Runtime comparison",
    "• If not executed: NOT_EXECUTED with reason, no invented numbers"
], images_dir/"threshold_performance.png")

# Slide 18 Benchmark
add_content_slide("Benchmark Results and Model Comparison", [
    f"• Model comparison: best {best_model} Val PR-AUC {best_score}",
    "• Traditional: DecisionTree/RF/HGB test PR-AUC 0.98-1.0 (synthetic separable, income+amount dominant)",
    f"• DL: test PR-AUC {dl_test} — underperforms trees",
    "• Anomaly: LOF best ROC 0.737 PR 0.147 — high FP",
    "• Document: bill total mismatch detection works (HIGH risk when mismatch)",
    "• Hybrid: sample FLAG_FOR_MANUAL_REVIEW due to doc FAILED even when ML prob 0.0 — conservative",
    "• Confusion matrix: TP 54 FP 1 FN 0 TN 845 @ thr 0.95",
    "• Feature importance: PatientIncome, ClaimAmount",
    "• Presentation and report use consistent results from evaluation/"
], images_dir/"confusion_matrix.png")

# Slide 19 Explainability
add_content_slide("Explainability, Risk Controls, and Human Review", [
    "• Feature importance: feature_importance.png, SHAP attempted (not installed)",
    "• Human explanation: key features, doc errors, policy violations, anomaly indicators, missing evidence, grounded, not vague",
    "• Risk controls: weights, thresholds, conservative review zone, doc FAILED => manual review, anomaly top_k => review, high amount => pre-auth check",
    "• Model cards: intended use decision-support, out-of-scope autonomous, training/eval data, metrics, limitations, bias, threshold, failure modes, version",
    "• Human review mandatory for REJECT_OR_ESCALATE, appeal mechanism",
    "• Disclaimer in every output: not final legal/insurance determination",
    "• Distinguish: confirmed label vs prediction vs anomaly vs rule violation vs missing evidence vs human decision"
], images_dir/"feature_importance.png")

# Slide 20 Limitations
add_content_slide("Limitations, Future Work, Conclusion, and References", [
    "• Limitations: synthetic 4500 modest, no free-text notes, no real doc images, high-card OHE memory heavy (7GB covariance fail), ClaimStatus potential leakage, near-perfect ML not realistic, anomaly low PR, DL underperforms, OCR fallback, RAG TFIDF less semantic, no graph collusion, no FHIR",
    "• Future: larger CMS data, time/group split, embeddings for codes, graph features provider-patient network, fairness audit, calibration cost-sensitive, encrypted secure API, Next.js frontend reviewer workflow, monitoring drift retraining",
    "• Conclusion: Built complete reproducible end-to-end academic prototype with 6 approaches, eval, visuals, docs, presentation, report, runnable without paid APIs, responsible disclaimer, human-in-loop",
    "• References: CMS public files, Kaggle healthcare fraud, sklearn, imbalanced-learn, SHAP, RAG Lewis et al., Model Cards Mitchell et al., EU Trustworthy AI, see documentation/references.md",
    "• Commands: pip install -r requirements.txt; python approaches/01_traditional_ml.py; ... ; python visualization_generator.py; python presentation/generate_presentation.py; python report/generate_report.py",
])

out_path = pres_dir / "medical_insurance_fraud_detection.pptx"
prs.save(str(out_path))
print(f"Saved presentation to {out_path}")
